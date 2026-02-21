"""
Script para criar a apresentação PPTX do Bloco 6 — IA aplicada à Ciência de Dados
Curso ABAR de Medições Inteligentes e Gestão Integrada | Fevereiro 2026

Uso: python create_slides.py
"""

import os
import sys

# ─── Verificar/instalar python-pptx ───────────────────────────────────────────
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree

# ─── Caminhos ─────────────────────────────────────────────────────────────────
BASE = r'C:\Users\raula\OneDrive\Documentos\Codigos e sistemas pessoais\AGENERSA\Cursos ABAR de Dados\Apresentação Raul'
OUTPUTS = r'C:\Users\raula\OneDrive\Documentos\Codigos e sistemas pessoais\AGENERSA\Cursos ABAR de Dados\outputs'
SCREENS_PRES = os.path.join(BASE, 'screenshots')
SCREENS_WEB = os.path.join(OUTPUTS, 'screenshots')
DIAGRAMAS = os.path.join(OUTPUTS, 'diagramas')

OUTPUT_PPTX = os.path.join(BASE, 'Bloco6_AI_Data_Science_ABAR.pptx')

# Imagens
IMG_IMO = os.path.join(SCREENS_PRES, 'benchmark_imo_2025.png')
IMG_SWEBENCH = os.path.join(SCREENS_PRES, 'benchmark_swebench.png')
IMG_EPOCH = os.path.join(SCREENS_PRES, 'benchmark_epoch_ai.png')
IMG_CHATBOT = os.path.join(SCREENS_PRES, 'benchmark_chatbot_arena.png')
IMG_ARC = os.path.join(SCREENS_PRES, 'benchmark_arc_agi.png')
IMG_STANFORD = os.path.join(SCREENS_PRES, 'benchmark_ai_timeline_stanford.png')
IMG_IBM = os.path.join(BASE, 'WhatsApp Image 2026-02-12 at 10.32.47.jpeg')
IMG_WEB1 = os.path.join(SCREENS_WEB, '01_web_config.png')
IMG_WEB2 = os.path.join(SCREENS_WEB, '02_web_pipeline.png')
IMG_WEB3 = os.path.join(SCREENS_WEB, '03_web_graficos.png')
IMG_WEB4 = os.path.join(SCREENS_WEB, '04_web_diagramas.png')
IMG_WEB5 = os.path.join(SCREENS_WEB, '05_web_textos.png')
IMG_WEB6 = os.path.join(SCREENS_WEB, '06_web_downloads.png')
IMG_ESTRUTURA = os.path.join(DIAGRAMAS, 'estrutura_distrito.png')
IMG_FLUXO = os.path.join(DIAGRAMAS, 'fluxo_auditoria.png')
IMG_PROCESSO = os.path.join(DIAGRAMAS, 'processo_analise.png')
IMG_LOGO = os.path.join(BASE, 'template_Imagem 10.jpg')
IMG_BG = os.path.join(BASE, 'template_bg.png')

# ─── Cores ────────────────────────────────────────────────────────────────────
COR_AZUL_ABAR      = RGBColor(0x00, 0x38, 0x86)   # #003886 - azul institucional ABAR
COR_AZUL_CLARO     = RGBColor(0x00, 0x81, 0xBF)   # #0081BF - azul claro do template
COR_BRANCO         = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
COR_CINZA_CLARO    = RGBColor(0xF0, 0xF4, 0xF8)   # cinza azulado suave
COR_CINZA_MEDIO    = RGBColor(0xD0, 0xD8, 0xE0)   # cinza médio para bordas
COR_CINZA_TEXTO    = RGBColor(0x44, 0x44, 0x55)   # cinza escuro para texto
COR_VERDE_DESTAQUE = RGBColor(0x00, 0x8A, 0x4B)   # verde para destaques
COR_LARANJA        = RGBColor(0xE8, 0x6D, 0x00)   # laranja para alertas/destaque
COR_AMARELO_DEMO   = RGBColor(0xFF, 0xD7, 0x00)   # amarelo para fundo demo

# ─── Dimensões (widescreen 16:9) ──────────────────────────────────────────────
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

# ─── Fontes ───────────────────────────────────────────────────────────────────
FONTE_TITULO = "Calibri"
FONTE_CORPO  = "Calibri"

# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_image_safe(slide, img_path, left, top, width, height):
    """Adiciona imagem se existir, senão retângulo placeholder."""
    if os.path.exists(img_path):
        return slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COR_CINZA_MEDIO
        shape.line.color.rgb = COR_AZUL_CLARO
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"[Imagem: {os.path.basename(img_path)}]"
        run.font.size = Pt(9)
        run.font.color.rgb = COR_CINZA_TEXTO
        return shape


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    """Adiciona retângulo com cor sólida."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=COR_CINZA_TEXTO, align=PP_ALIGN.LEFT,
                font_name=FONTE_CORPO, word_wrap=True):
    """Adiciona caixa de texto simples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = word_wrap
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return txBox


def add_paragraph(text_frame, text, font_size=12, bold=False,
                  color=COR_CINZA_TEXTO, align=PP_ALIGN.LEFT,
                  font_name=FONTE_CORPO, space_before=0, indent_level=0):
    """Adiciona parágrafo num text_frame existente."""
    p = text_frame.add_paragraph()
    p.alignment = align
    p.level = indent_level
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return p


def setup_slide_background(slide, color=COR_BRANCO):
    """Define cor de fundo do slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text, subtitle_text=None,
                   bar_color=COR_AZUL_ABAR, text_color=COR_BRANCO,
                   bar_height_cm=3.0):
    """Adiciona barra de cabeçalho azul com título."""
    bar_h = Cm(bar_height_cm)
    bar = add_rect(slide, Cm(0), Cm(0), SLIDE_W, bar_h, bar_color)

    # Título na barra
    title_box = slide.shapes.add_textbox(
        Cm(0.8), Cm(0.2),
        SLIDE_W - Cm(1.6), bar_h - Cm(0.4)
    )
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24) if subtitle_text else Pt(26)
    run.font.bold = True
    run.font.name = FONTE_TITULO
    run.font.color.rgb = text_color

    if subtitle_text:
        sub = add_textbox(
            slide,
            Cm(0.8), bar_h - Cm(0.05),
            SLIDE_W - Cm(1.6), Cm(0.9),
            subtitle_text,
            font_size=13,
            bold=False,
            color=COR_AZUL_CLARO,
            align=PP_ALIGN.LEFT
        )

    return bar_h


def add_footer(slide, text="Curso ABAR — Medições Inteligentes e Gestão Integrada | Fevereiro 2026"):
    """Adiciona rodapé no slide."""
    footer_top = SLIDE_H - Cm(0.7)
    add_rect(slide, Cm(0), footer_top, SLIDE_W, Cm(0.7), COR_AZUL_ABAR)
    add_textbox(
        slide,
        Cm(0.5), footer_top + Cm(0.05),
        SLIDE_W - Cm(1), Cm(0.6),
        text,
        font_size=8,
        color=COR_BRANCO,
        align=PP_ALIGN.CENTER
    )


def add_logo(slide):
    """Adiciona logo ABAR no canto superior direito."""
    logo_w = Cm(5.0)
    logo_h = Cm(1.5)
    logo_left = SLIDE_W - logo_w - Cm(0.3)
    logo_top = Cm(0.2)
    if os.path.exists(IMG_LOGO):
        slide.shapes.add_picture(IMG_LOGO, logo_left, logo_top, logo_w, logo_h)


def add_demo_banner(slide, top_offset=Cm(3.2)):
    """Adiciona banner amarelo DEMO AO VIVO."""
    banner = add_rect(slide, Cm(0.8), top_offset, Cm(10), Cm(1.0), COR_AMARELO_DEMO, COR_LARANJA, Pt(2))
    add_textbox(
        slide,
        Cm(0.8), top_offset,
        Cm(10), Cm(1.0),
        "★ DEMO AO VIVO",
        font_size=14,
        bold=True,
        color=RGBColor(0x80, 0x30, 0x00),
        align=PP_ALIGN.CENTER
    )


def add_card(slide, left, top, width, height, titulo, corpo,
             card_color=COR_CINZA_CLARO, title_color=COR_AZUL_ABAR):
    """Adiciona card com título e corpo."""
    # Sombra
    shadow = add_rect(slide, left + Cm(0.15), top + Cm(0.15),
                      width, height, COR_CINZA_MEDIO)
    # Card
    card = add_rect(slide, left, top, width, height, card_color, COR_AZUL_CLARO, Pt(1))
    # Título
    title_bg = add_rect(slide, left, top, width, Cm(0.9), title_color)
    add_textbox(slide, left + Cm(0.2), top + Cm(0.05),
                width - Cm(0.4), Cm(0.8),
                titulo, font_size=11, bold=True,
                color=COR_BRANCO, align=PP_ALIGN.LEFT)
    # Corpo
    body_box = slide.shapes.add_textbox(
        left + Cm(0.2), top + Cm(1.0),
        width - Cm(0.4), height - Cm(1.1)
    )
    body_box.text_frame.word_wrap = True
    p = body_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = corpo
    run.font.size = Pt(10)
    run.font.name = FONTE_CORPO
    run.font.color.rgb = COR_CINZA_TEXTO
    return card


# ─── Criação da Apresentação ─────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # Blank layout

print("Criando 25 slides...")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_AZUL_ABAR)

# Faixa decorativa esquerda
add_rect(slide, Cm(0), Cm(0), Cm(0.6), SLIDE_H, COR_AZUL_CLARO)

# Painel branco central
add_rect(slide, Cm(0.6), Cm(2.5), SLIDE_W - Cm(0.6) - Cm(3.5), SLIDE_H - Cm(4.5), COR_BRANCO)

# Logo ABAR
if os.path.exists(IMG_LOGO):
    slide.shapes.add_picture(IMG_LOGO, Cm(1.2), Cm(0.4), Cm(6.5), Cm(2.0))

# Linha decorativa abaixo do logo
add_rect(slide, Cm(0.6), Cm(2.5), SLIDE_W - Cm(0.6) - Cm(3.5), Cm(0.12), COR_AZUL_CLARO)

# Bloco do título
title_box = slide.shapes.add_textbox(Cm(1.2), Cm(3.0), Cm(25), Cm(4.0))
title_box.text_frame.word_wrap = True
p = title_box.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Bloco 6 — Auditoria de Dados,"
run.font.size = Pt(28)
run.font.bold = True
run.font.name = FONTE_TITULO
run.font.color.rgb = COR_AZUL_ABAR

add_paragraph(title_box.text_frame, "BI e Transparência Reguladora",
              font_size=28, bold=True, color=COR_AZUL_ABAR, align=PP_ALIGN.LEFT)

# Subtítulo
sub_box = slide.shapes.add_textbox(Cm(1.2), Cm(7.3), Cm(25), Cm(2.0))
sub_box.text_frame.word_wrap = True
p = sub_box.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Como a Inteligência Artificial transforma"
run.font.size = Pt(17)
run.font.name = FONTE_CORPO
run.font.color.rgb = COR_AZUL_CLARO

add_paragraph(sub_box.text_frame, "a ciência de dados para reguladores",
              font_size=17, color=COR_AZUL_CLARO, align=PP_ALIGN.LEFT)

# Linha divisória
add_rect(slide, Cm(1.2), Cm(9.6), Cm(22), Cm(0.06), COR_AZUL_CLARO)

# Nome do professor
add_textbox(slide, Cm(1.2), Cm(10.0), Cm(22), Cm(0.9),
            "Prof. Raul Araújo — AGENERSA",
            font_size=15, bold=True, color=COR_AZUL_ABAR)

# Rodapé da capa
add_rect(slide, Cm(0.6), SLIDE_H - Cm(1.5), SLIDE_W - Cm(0.6) - Cm(3.5), Cm(1.5), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.2), SLIDE_H - Cm(1.4), Cm(22), Cm(1.2),
            "Curso ABAR — Medições Inteligentes e Gestão Integrada  |  Fevereiro 2026",
            font_size=11, color=COR_BRANCO, align=PP_ALIGN.LEFT)

# Imagem decorativa (background lateral)
add_image_safe(slide, IMG_BG,
               SLIDE_W - Cm(8.5), Cm(0),
               Cm(8.5), SLIDE_H)

print("  Slide 1 OK - Capa")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — A IA que mudou tudo em 2 anos
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "A IA que Mudou Tudo em 2 Anos",
                                 subtitle_text="De curiosidade acadêmica a ferramenta profissional indispensável")
add_logo(slide)
add_footer(slide)

# Coluna esquerda — bullets
bullets = [
    ("2022", "IA gerava texto confuso, errava contas simples — considerada brinquedo"),
    ("2024", "GPT-4 aprovado no Bar Exam (top 10%) e no LSAT — nível de advogado"),
    ("2025", "Gemini Deep Think → Medalha de Ouro na IMO (Olimpíada Internacional de Matemática)"),
    ("2025", "Claude 4.5 Opus → resolve 76,8% de bugs reais em repositórios GitHub (SWE-bench)"),
    ("Reflexão", "\"Se uma IA resolve olimpíada de matemática... o que pode fazer com seus dados?\""),
]

y = content_bottom + Cm(0.4)
for ano, texto in bullets:
    # Badge do ano
    badge_color = COR_AZUL_ABAR if ano != "Reflexão" else COR_VERDE_DESTAQUE
    badge = add_rect(slide, Cm(0.8), y, Cm(2.0), Cm(0.75), badge_color)
    add_textbox(slide, Cm(0.8), y, Cm(2.0), Cm(0.75),
                ano, font_size=10, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    # Texto
    txt_box = slide.shapes.add_textbox(Cm(3.2), y, Cm(15.8), Cm(0.85))
    txt_box.text_frame.word_wrap = True
    p = txt_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(12.5)
    run.font.name = FONTE_CORPO
    run.font.color.rgb = COR_AZUL_ABAR if ano == "Reflexão" else COR_CINZA_TEXTO
    run.font.bold = (ano == "Reflexão")
    run.font.italic = (ano == "Reflexão")
    y += Cm(1.35)

# Imagem lateral
add_image_safe(slide, IMG_IMO,
               Cm(20.5), content_bottom + Cm(0.3),
               Cm(12.6), Cm(13.5))

print("  Slide 2 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Roteiro das 4 Horas
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Roteiro das 4 Horas",
                                 subtitle_text="Fio condutor: IBM Data Science Methodology aplicada com IA")
add_logo(slide)
add_footer(slide)

modulos = [
    ("Módulo 1", "30 min", "Ecossistema de ferramentas IA\nChatGPT, Gemini, Copilot, Claude Code\nNíveis: do dia a dia ao projeto completo", COR_AZUL_ABAR),
    ("Módulo 2", "60 min", "Ciclo IBM Data Science com IA\nCada etapa acelerada por IA\nDe dados brutos ao modelo em produção", COR_AZUL_CLARO),
    ("Módulo 3", "50 min", "Demo ao vivo — do Excel ao relatório\n7 Notebooks, 28 chamadas Gemini\nRelatório de 8,9 MB em 1 minuto e 42 s", COR_VERDE_DESTAQUE),
    ("Módulo 4", "15 min", "NotebookLM para os alunos\nCarregar apostila do Prof. Barateiro\nPerguntas técnicas em linguagem natural", RGBColor(0x8B, 0x00, 0x8B)),
]

block_w = Cm(7.5)
block_h = Cm(11.5)
x_start = Cm(0.8)
y_start = content_bottom + Cm(0.5)
gap = Cm(0.55)

for i, (titulo, duracao, corpo, cor) in enumerate(modulos):
    x = x_start + i * (block_w + gap)
    # Bloco principal
    add_rect(slide, x, y_start, block_w, block_h, COR_CINZA_CLARO, cor, Pt(2))
    # Header do bloco
    add_rect(slide, x, y_start, block_w, Cm(1.8), cor)
    add_textbox(slide, x + Cm(0.2), y_start + Cm(0.1),
                block_w - Cm(0.4), Cm(0.8),
                titulo, font_size=14, bold=True,
                color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Cm(0.2), y_start + Cm(0.95),
                block_w - Cm(0.4), Cm(0.75),
                duracao, font_size=18, bold=True,
                color=COR_AMARELO_DEMO, align=PP_ALIGN.CENTER)
    # Corpo
    linhas = corpo.split('\n')
    y_txt = y_start + Cm(2.1)
    for linha in linhas:
        txt_b = slide.shapes.add_textbox(
            x + Cm(0.3), y_txt,
            block_w - Cm(0.6), Cm(0.85)
        )
        txt_b.text_frame.word_wrap = True
        p = txt_b.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "• " + linha
        run.font.size = Pt(10.5)
        run.font.name = FONTE_CORPO
        run.font.color.rgb = COR_CINZA_TEXTO
        y_txt += Cm(1.0)

# Nota intervalo
intervalo_box = add_rect(slide,
                          Cm(0.8), y_start + block_h + Cm(0.3),
                          Cm(31.5), Cm(0.7),
                          COR_CINZA_MEDIO)
add_textbox(slide,
            Cm(0.8), y_start + block_h + Cm(0.32),
            Cm(31.5), Cm(0.65),
            "⏸  Intervalo de 10 minutos entre Módulos 2 e 3",
            font_size=11, bold=False, color=COR_AZUL_ABAR, align=PP_ALIGN.CENTER)

print("  Slide 3 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Ferramentas Nível 1: Dia a Dia
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Nível 1: Ferramentas para o Dia a Dia",
                                 subtitle_text="Para quem não programa e não quer programar")
add_logo(slide)
add_footer(slide)

ferramentas = [
    ("ChatGPT na Web", "chatgpt.com",
     "Análise de planilhas diretamente pelo navegador\nGeração de gráficos automáticos\nResumos de documentos longos\nResposta em segundos",
     COR_AZUL_ABAR),
    ("Google Gemini", "gemini.google.com",
     "Integrado ao Google Drive e Docs\nAnálise de Sheets sem sair do Google\nMultimodal: texto, imagem e áudio\nGratuito com conta Google",
     COR_AZUL_CLARO),
    ("Microsoft Copilot", "copilot.microsoft.com",
     "Integrado ao Office 365 (Excel, Word, PPT)\nAnálise de dados no próprio Excel\nGeração de apresentações\nDisponível em ambientes corporativos",
     COR_VERDE_DESTAQUE),
]

card_w = Cm(10.0)
card_h = Cm(10.8)
x_start = Cm(0.8)
y_start = content_bottom + Cm(0.5)
gap = Cm(0.65)

for i, (nome, url, corpo, cor) in enumerate(ferramentas):
    x = x_start + i * (card_w + gap)
    add_rect(slide, x, y_start, card_w, card_h, COR_CINZA_CLARO, cor, Pt(2))
    add_rect(slide, x, y_start, card_w, Cm(1.6), cor)
    add_textbox(slide, x + Cm(0.2), y_start + Cm(0.1),
                card_w - Cm(0.4), Cm(0.8),
                nome, font_size=13, bold=True,
                color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Cm(0.2), y_start + Cm(0.9),
                card_w - Cm(0.4), Cm(0.6),
                url, font_size=9, bold=False,
                color=COR_AMARELO_DEMO, align=PP_ALIGN.CENTER)
    linhas = corpo.split('\n')
    y_txt = y_start + Cm(1.8)
    for linha in linhas:
        tb = slide.shapes.add_textbox(x + Cm(0.3), y_txt, card_w - Cm(0.6), Cm(0.85))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "✓ " + linha
        run.font.size = Pt(10.5)
        run.font.name = FONTE_CORPO
        run.font.color.rgb = COR_CINZA_TEXTO
        y_txt += Cm(1.0)

# Destaque caso de uso
dest_y = y_start + card_h + Cm(0.4)
add_rect(slide, Cm(0.8), dest_y, Cm(32.0), Cm(1.0), COR_AZUL_ABAR)
add_textbox(slide, Cm(0.8), dest_y + Cm(0.1),
            Cm(32.0), Cm(0.85),
            "CASO DE USO:  Cole uma tabela → peça análise → receba gráfico em 30 segundos",
            font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 4 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Ferramentas Nível 2: Automação
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Nível 2: Ferramentas para Automação",
                                 subtitle_text="Para quem quer automatizar sem saber programar de verdade")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

ferramentas2 = [
    ("VS Code + GitHub Copilot", "Editor de código gratuito com IA integrada. Sugere código linha a linha enquanto você escreve em linguagem natural."),
    ("Cursor AI", "IDE focada em IA. Toda a interface foi redesenhada para geração de código por conversa. O mais popular entre analistas de dados que estão começando."),
    ("JetBrains AI", "Suite profissional (PyCharm, DataSpell) com IA para Python e análise de dados. Ideal para ambientes corporativos."),
]

for nome, desc in ferramentas2:
    bg = add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(2.4), COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(1))
    add_rect(slide, Cm(0.8), y, Cm(0.4), Cm(2.4), COR_AZUL_CLARO)
    add_textbox(slide, Cm(1.5), y + Cm(0.1), Cm(31.0), Cm(0.8),
                nome, font_size=13, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, Cm(1.5), y + Cm(0.9), Cm(31.0), Cm(1.3),
                desc, font_size=11, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(2.7)

# Nota de aprendizado
add_rect(slide, Cm(0.8), y + Cm(0.2), Cm(32.0), Cm(1.8), RGBColor(0xFF, 0xF8, 0xE1), COR_LARANJA, Pt(1))
add_textbox(slide, Cm(1.2), y + Cm(0.3), Cm(31.5), Cm(1.5),
            "Diferencial: Você descreve o que quer em português → IA escreve o código em Python\n"
            "Aprendizado implícito: você absorve programação naturalmente vendo a IA trabalhar",
            font_size=12, bold=False, color=COR_CINZA_TEXTO, word_wrap=True)

print("  Slide 5 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Ferramentas Nível 3: Projetos Completos
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Nível 3: IA para Projetos Completos",
                                 subtitle_text="Para pipelines de análise de dados end-to-end")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

ferramentas3 = [
    ("Claude Code CLI", "Agente de IA autônomo em linha de comando. Lê e escreve arquivos, executa código, cria sistemas completos a partir de uma descrição em texto."),
    ("Aider", "Assistente de programação em terminal. Conecta ao seu repositório Git e implementa funcionalidades completas com um único comando."),
    ("Devin (Cognition AI)", "Primeiro agente de software autônomo. Planeja, pesquisa, codifica, testa e implanta — sem supervisão constante."),
]

for nome, desc in ferramentas3:
    bg = add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(2.4), COR_CINZA_CLARO, COR_AZUL_ABAR, Pt(1))
    add_rect(slide, Cm(0.8), y, Cm(0.4), Cm(2.4), COR_AZUL_ABAR)
    add_textbox(slide, Cm(1.5), y + Cm(0.1), Cm(31.0), Cm(0.8),
                nome, font_size=13, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, Cm(1.5), y + Cm(0.9), Cm(31.0), Cm(1.3),
                desc, font_size=11, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(2.7)

# Destaque especial
add_rect(slide, Cm(0.8), y + Cm(0.2), Cm(32.0), Cm(2.0), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.2), y + Cm(0.3), Cm(31.5), Cm(1.6),
            "★  ESTE PROJETO foi construído inteiramente com Claude Code\n"
            "   7 notebooks Python, interface web FastAPI, relatório Word de 100+ páginas — zero linha de código manual",
            font_size=12, bold=True, color=COR_BRANCO, word_wrap=True)

print("  Slide 6 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Escolhendo a Ferramenta Certa
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Escolhendo a Ferramenta Certa",
                                 subtitle_text="A ferramenta certa para cada nível de complexidade")
add_logo(slide)
add_footer(slide)

# Tabela
col_widths = [Cm(9.5), Cm(11.0), Cm(11.8)]
col_labels = ["Cenário", "Ferramenta Recomendada", "Exemplo Prático"]
rows = [
    ("Análise rápida, dado pronto", "ChatGPT web / Gemini", "Colar planilha e pedir gráfico de série temporal"),
    ("Script repetitivo, simples", "GitHub Copilot / Cursor AI", "Criar rotina de limpeza de dados em Excel"),
    ("Pipeline completo, projeto", "Claude Code / Aider", "Sistema de auditoria automatizada com relatório"),
    ("Documentos técnicos", "NotebookLM", "Consultar normas GUM e apostila do curso"),
    ("Imagens e apresentações", "ChatGPT / Copilot", "Gerar diagrama a partir de texto descritivo"),
]

table_x = Cm(0.8)
table_y = content_bottom + Cm(0.5)
row_h = Cm(1.7)
total_w = sum(col_widths)

# Cabeçalho
x = table_x
for j, (label, w) in enumerate(zip(col_labels, col_widths)):
    add_rect(slide, x, table_y, w, row_h, COR_AZUL_ABAR, COR_BRANCO, Pt(1))
    add_textbox(slide, x + Cm(0.2), table_y + Cm(0.2), w - Cm(0.4), row_h - Cm(0.4),
                label, font_size=12, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    x += w

for i, row in enumerate(rows):
    y_row = table_y + (i + 1) * row_h
    x = table_x
    bg_color = COR_CINZA_CLARO if i % 2 == 0 else COR_BRANCO
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        add_rect(slide, x, y_row, w, row_h, bg_color, COR_CINZA_MEDIO, Pt(1))
        text_color = COR_AZUL_ABAR if j == 1 else COR_CINZA_TEXTO
        bold = (j == 1)
        add_textbox(slide, x + Cm(0.2), y_row + Cm(0.15), w - Cm(0.4), row_h - Cm(0.3),
                    cell, font_size=10.5, bold=bold, color=text_color,
                    align=PP_ALIGN.LEFT, word_wrap=True)
        x += w

# Mensagem chave
bottom_y = table_y + (len(rows) + 1) * row_h + Cm(0.4)
add_rect(slide, Cm(0.8), bottom_y, total_w, Cm(1.0), COR_VERDE_DESTAQUE)
add_textbox(slide, Cm(1.0), bottom_y + Cm(0.1), total_w - Cm(0.4), Cm(0.85),
            "Princípio: comece com a ferramenta mais simples que resolve o seu problema",
            font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 7 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEMO 1: ChatGPT + Planilha
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, RGBColor(0xFF, 0xFF, 0xF0))
content_bottom = add_header_bar(slide, "DEMO AO VIVO — ChatGPT + Planilha Excel",
                                 subtitle_text="Módulo 1 • Demo 1 de 4",
                                 bar_color=COR_LARANJA)
add_logo(slide)
add_footer(slide)
add_demo_banner(slide, content_bottom + Cm(0.3))

y = content_bottom + Cm(1.6)

passos = [
    ("1", "Acessar chatgpt.com", "Abrir o navegador e fazer login na conta ChatGPT (gratuita)"),
    ("2", "Fazer upload da planilha", "Clicar no ícone de clipe → selecionar aba da planilha do Distrito de Gás"),
    ("3", "Digitar o pedido", "\"Faça uma análise exploratória desta tabela e gere um gráfico de série temporal do volume diário\""),
    ("4", "Aguardar e interpretar", "Em ~30 segundos: estatísticas descritivas + gráfico interativo + insights automáticos"),
]

for num, titulo, desc in passos:
    # Círculo numerado
    circ = add_rect(slide, Cm(0.8), y, Cm(1.2), Cm(1.2), COR_LARANJA)
    add_textbox(slide, Cm(0.8), y, Cm(1.2), Cm(1.2),
                num, font_size=16, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(2.3), y, Cm(30.0), Cm(0.65),
                titulo, font_size=13, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, Cm(2.3), y + Cm(0.65), Cm(30.0), Cm(0.75),
                desc, font_size=11, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(1.85)

# Nota de rodapé informativa
add_rect(slide, Cm(0.8), y + Cm(0.3), Cm(32.0), Cm(1.0), COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(1))
add_textbox(slide, Cm(1.2), y + Cm(0.4), Cm(31.5), Cm(0.8),
            "Dados utilizados: Planilha Distrito de Gás — 183 dias de medições (14 abas, clientes residenciais e industriais)",
            font_size=11, color=COR_AZUL_ABAR, word_wrap=True)

print("  Slide 8 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — IBM Data Science Methodology + IA
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "IBM Data Science Methodology + IA",
                                 subtitle_text="O ciclo completo de ciência de dados — acelerado por IA em cada etapa")
add_logo(slide)
add_footer(slide)

# Imagem do diagrama IBM — central
img_top = content_bottom + Cm(0.4)
img_h = SLIDE_H - img_top - Cm(1.0)
img_w = Cm(13.0)
img_left = (SLIDE_W - img_w) / 2

add_image_safe(slide, IMG_IBM, img_left, img_top, img_w, img_h)

# Anotações ao redor
anotacoes_esq = [
    (Cm(0.5), img_top + Cm(0.5), "Entendimento\ndo Negócio", "NotebookLM\n+ apostila"),
    (Cm(0.5), img_top + Cm(3.5), "Coleta\ndos Dados", "Web scraping\ncom Claude"),
    (Cm(0.5), img_top + Cm(6.5), "Preparação\ndos Dados", "IA gera código\nPython"),
]

anotacoes_dir = [
    (SLIDE_W - Cm(9.5), img_top + Cm(0.5), "Criação\ndo Modelo", "IA lê metodologia\ntécnica GUM"),
    (SLIDE_W - Cm(9.5), img_top + Cm(3.5), "Avaliação", "Validação\nautomática"),
    (SLIDE_W - Cm(9.5), img_top + Cm(6.5), "Deployment", "Pipeline FastAPI\nautomático"),
]

for x, y, etapa, ferr in anotacoes_esq:
    add_rect(slide, x, y, Cm(7.5), Cm(2.5), COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(1))
    add_textbox(slide, x + Cm(0.2), y + Cm(0.1), Cm(7.1), Cm(1.0),
                etapa, font_size=11, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, x + Cm(0.2), y + Cm(1.1), Cm(7.1), Cm(1.2),
                ferr, font_size=10, color=COR_CINZA_TEXTO, word_wrap=True)

for x, y, etapa, ferr in anotacoes_dir:
    add_rect(slide, x, y, Cm(8.5), Cm(2.5), COR_CINZA_CLARO, COR_AZUL_ABAR, Pt(1))
    add_textbox(slide, x + Cm(0.2), y + Cm(0.1), Cm(8.1), Cm(1.0),
                etapa, font_size=11, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, x + Cm(0.2), y + Cm(1.1), Cm(8.1), Cm(1.2),
                ferr, font_size=10, color=COR_CINZA_TEXTO, word_wrap=True)

print("  Slide 9 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Etapa 1: Entendimento do Negócio
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Etapa 1: Entendimento do Negócio",
                                 subtitle_text="IBM Data Science Methodology • Passo 1 de 7")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

# Problema
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(2.2), COR_CINZA_CLARO, COR_AZUL_ABAR, Pt(1))
add_textbox(slide, Cm(1.1), y + Cm(0.1), Cm(31.5), Cm(0.75),
            "DESAFIO:", font_size=12, bold=True, color=COR_AZUL_ABAR)
add_textbox(slide, Cm(1.1), y + Cm(0.85), Cm(31.5), Cm(1.1),
            "Como entender a metodologia GUM de incertezas sem ler 200 páginas de norma técnica?",
            font_size=13, color=COR_CINZA_TEXTO, word_wrap=True)

y += Cm(2.7)

# Solução
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(0.7), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.05), Cm(31.5), Cm(0.65),
            "SOLUÇÃO: NotebookLM — carregue a apostila e faça perguntas",
            font_size=13, bold=True, color=COR_BRANCO)

y += Cm(1.1)

# Passos
passos_notebooklm = [
    "Acesse notebooklm.google.com (gratuito com conta Google)",
    "Carregue a apostila do Prof. Barateiro como fonte",
    "Faça perguntas técnicas em linguagem natural",
    "NotebookLM responde com citação exata da página",
]

for passo in passos_notebooklm:
    add_rect(slide, Cm(0.8), y, Cm(0.6), Cm(0.8), COR_AZUL_CLARO)
    add_textbox(slide, Cm(1.7), y, Cm(31.0), Cm(0.85),
                passo, font_size=12, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(1.05)

y += Cm(0.2)

# Exemplos de perguntas
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(0.65), COR_CINZA_MEDIO)
add_textbox(slide, Cm(1.0), y + Cm(0.05), Cm(31.5), Cm(0.6),
            "Exemplo de pergunta ao vivo:", font_size=12, bold=True, color=COR_AZUL_ABAR)

y += Cm(0.85)
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(1.5), RGBColor(0xE8, 0xF4, 0xFD), COR_AZUL_CLARO, Pt(1))
add_textbox(slide, Cm(1.2), y + Cm(0.1), Cm(31.2), Cm(1.3),
            "\"Qual é a equação de incerteza combinada do tipo A segundo a norma GUM?\"\n"
            "\"Quais são os requisitos de medição fiscal em dutos offshore?\"",
            font_size=12, color=COR_AZUL_ABAR, word_wrap=True)

y += Cm(1.8)
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(0.9), COR_VERDE_DESTAQUE)
add_textbox(slide, Cm(1.0), y + Cm(0.1), Cm(31.5), Cm(0.75),
            "Resultado: o regulador consulta em segundos, sem precisar ler o documento completo",
            font_size=12, bold=True, color=COR_BRANCO, word_wrap=True)

print("  Slide 10 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Etapa 2: Coleta de Dados
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Etapa 2: Coleta de Dados",
                                 subtitle_text="IBM Data Science Methodology • Passo 2 de 7")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

# Cenário
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(1.8), COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(1))
add_textbox(slide, Cm(1.0), y + Cm(0.1), Cm(31.5), Cm(0.7),
            "CENÁRIO REAL:", font_size=12, bold=True, color=COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.8), Cm(31.5), Cm(0.85),
            "Dados da ANP, ANEEL e distribuidoras — espalhados em dezenas de páginas na web",
            font_size=12, color=COR_CINZA_TEXTO, word_wrap=True)

y += Cm(2.3)

# Comparativo Antes x Depois
col_w = Cm(15.5)
# Antes
add_rect(slide, Cm(0.8), y, col_w, Cm(5.5), RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xCC, 0x00, 0x00), Pt(1))
add_rect(slide, Cm(0.8), y, col_w, Cm(0.9), RGBColor(0xCC, 0x00, 0x00))
add_textbox(slide, Cm(0.9), y + Cm(0.05), col_w - Cm(0.2), Cm(0.8),
            "ANTES (sem IA)", font_size=12, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

antes = [
    "Copiar/colar manualmente de cada página",
    "Horas ou dias de trabalho repetitivo",
    "Erros de transcrição frequentes",
    "Difícil de automatizar e repetir",
]
y_a = y + Cm(1.1)
for item in antes:
    add_textbox(slide, Cm(1.1), y_a, col_w - Cm(0.6), Cm(0.75),
                "✗ " + item, font_size=11, color=RGBColor(0x80, 0x00, 0x00), word_wrap=True)
    y_a += Cm(0.85)

# Depois
x_dep = col_w + Cm(1.2)
add_rect(slide, x_dep, y, col_w, Cm(5.5), RGBColor(0xF0, 0xFF, 0xF4), COR_VERDE_DESTAQUE, Pt(1))
add_rect(slide, x_dep, y, col_w, Cm(0.9), COR_VERDE_DESTAQUE)
add_textbox(slide, x_dep + Cm(0.1), y + Cm(0.05), col_w - Cm(0.2), Cm(0.8),
            "DEPOIS (com Claude Code)", font_size=12, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

depois = [
    '"Me colete dados mensais de gás da ANP dos últimos 12 meses"',
    "Claude escreve o script → executa → retorna CSV organizado",
    "Tempo: minutos, não dias",
    "Script reutilizável para execuções futuras",
]
y_d = y + Cm(1.1)
for item in depois:
    add_textbox(slide, x_dep + Cm(0.2), y_d, col_w - Cm(0.5), Cm(0.75),
                "✓ " + item, font_size=11, color=RGBColor(0x00, 0x60, 0x30), word_wrap=True)
    y_d += Cm(0.85)

# Nota técnica
y_nota = y + Cm(5.8)
add_rect(slide, Cm(0.8), y_nota, Cm(32.0), Cm(1.0), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y_nota + Cm(0.1), Cm(31.5), Cm(0.85),
            "Técnica: Firecrawl API + Claude Code = web scraping inteligente com compreensão de contexto",
            font_size=12, bold=True, color=COR_BRANCO, word_wrap=True)

print("  Slide 11 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Limpeza e Preparação
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Etapas 3-4: Limpeza e Preparação dos Dados",
                                 subtitle_text="IBM Data Science Methodology • Passos 3-4 de 7")
add_logo(slide)
add_footer(slide)

# Imagem lateral
add_image_safe(slide, IMG_EPOCH,
               Cm(21.0), content_bottom + Cm(0.5),
               Cm(12.0), Cm(10.0))

y = content_bottom + Cm(0.5)

# Problema real
add_rect(slide, Cm(0.8), y, Cm(19.5), Cm(2.0), RGBColor(0xFF, 0xF5, 0xE0), COR_LARANJA, Pt(1))
add_textbox(slide, Cm(1.0), y + Cm(0.1), Cm(19.0), Cm(0.65),
            "PROBLEMA REAL:", font_size=12, bold=True, color=COR_LARANJA)
add_textbox(slide, Cm(1.0), y + Cm(0.75), Cm(19.0), Cm(1.1),
            "Planilha Excel com 14 abas, 183 dias, dados inconsistentes em unidades e formatação",
            font_size=11, color=COR_CINZA_TEXTO, word_wrap=True)

y += Cm(2.5)

# Prompt Claude
add_rect(slide, Cm(0.8), y, Cm(19.5), Cm(0.65), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.05), Cm(19.0), Cm(0.6),
            "PROMPT PARA O CLAUDE CODE:", font_size=11, bold=True, color=COR_BRANCO)
y += Cm(0.75)
add_rect(slide, Cm(0.8), y, Cm(19.5), Cm(2.8), RGBColor(0xF5, 0xF5, 0xFF), COR_AZUL_CLARO, Pt(1))
add_textbox(slide, Cm(1.1), y + Cm(0.1), Cm(19.0), Cm(2.6),
            '"Leia a planilha Excel de medições do Distrito, identifique valores '
            'ausentes, normalize as unidades (Mm³ e kcal/m³), remova outliers '
            'acima de 3 desvios padrão, e crie um DataFrame limpo exportado como CSV."',
            font_size=11, color=COR_AZUL_ABAR, word_wrap=True)

y += Cm(3.2)

# Resultado
resultados = [
    "Código Python gerado automaticamente",
    "183 dias de dados limpos e consistentes",
    "Unidades normalizadas para análise",
    "Valores ausentes imputados com média móvel",
]

for res in resultados:
    add_rect(slide, Cm(0.8), y, Cm(0.5), Cm(0.75), COR_VERDE_DESTAQUE)
    add_textbox(slide, Cm(1.6), y, Cm(18.5), Cm(0.8),
                res, font_size=11, color=COR_CINZA_TEXTO)
    y += Cm(0.95)

print("  Slide 12 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Criação do Modelo: IA lê a Apostila
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Etapa 5: Criação do Modelo — IA como Especialista Técnico",
                                 subtitle_text="IBM Data Science Methodology • Passo 5 de 7")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

# Inovação central
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(1.8), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.15), Cm(31.5), Cm(1.5),
            "INOVAÇÃO CENTRAL: E se a IA pudesse ler a apostila do Prof. Barateiro e aplicar a metodologia automaticamente?",
            font_size=14, bold=True, color=COR_BRANCO, word_wrap=True)

y += Cm(2.3)

# Fluxo
fluxo = [
    (COR_AZUL_ABAR, "Apostila PDF", "Arquivo do curso\ncomo fonte"),
    (COR_AZUL_CLARO, "Gemini extrai\nequações", "gemini-2.5-flash\nlê e interpreta"),
    (COR_VERDE_DESTAQUE, "Python\nimplementa", "Código gerado\nautomaticamente"),
    (COR_LARANJA, "Resultados\nvalidados", "183 dias\ncalculados"),
]

block_w = Cm(7.2)
block_h = Cm(3.5)
x = Cm(0.8)
gap = Cm(0.65)

for i, (cor, titulo, desc) in enumerate(fluxo):
    add_rect(slide, x, y, block_w, block_h, COR_CINZA_CLARO, cor, Pt(2))
    add_rect(slide, x, y, block_w, Cm(0.8), cor)
    add_textbox(slide, x + Cm(0.1), y + Cm(0.05),
                block_w - Cm(0.2), Cm(0.75),
                titulo, font_size=11, bold=True,
                color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Cm(0.2), y + Cm(1.0),
                block_w - Cm(0.4), Cm(2.3),
                desc, font_size=11, color=COR_CINZA_TEXTO,
                align=PP_ALIGN.CENTER, word_wrap=True)
    # Seta
    if i < len(fluxo) - 1:
        seta_x = x + block_w + Cm(0.1)
        add_textbox(slide, seta_x, y + block_h / 2 - Cm(0.3),
                    Cm(0.5), Cm(0.7), "→", font_size=20, bold=True,
                    color=COR_AZUL_ABAR, align=PP_ALIGN.CENTER)
    x += block_w + gap

y += block_h + Cm(0.8)

# Detalhes do arquivo
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(1.0), COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(1))
add_textbox(slide, Cm(1.0), y + Cm(0.1), Cm(31.5), Cm(0.85),
            "Arquivo: extrair_metodologia.py — extrai 7 seções de metodologia da apostila em formato Markdown",
            font_size=12, color=COR_AZUL_ABAR, word_wrap=True)

y += Cm(1.5)
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(1.0), COR_VERDE_DESTAQUE)
add_textbox(slide, Cm(1.0), y + Cm(0.1), Cm(31.5), Cm(0.85),
            "★  A metodologia GUM foi extraída automaticamente da apostila — sem digitar uma única equação",
            font_size=13, bold=True, color=COR_BRANCO, word_wrap=True)

print("  Slide 13 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Análise de Incertezas com IA
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Notebook 6: Incertezas de Medição (GUM) com IA",
                                 subtitle_text="Propagação de incerteza para 183 dias de dados de gás")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

# Equação
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(0.65), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.05), Cm(31.5), Cm(0.6),
            "Equação de Incerteza Combinada (GUM — extraída automaticamente da apostila):",
            font_size=12, bold=True, color=COR_BRANCO)

y += Cm(0.85)
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(2.0), RGBColor(0xF8, 0xF8, 0xFF), COR_AZUL_CLARO, Pt(1))
add_textbox(slide, Cm(2.0), y + Cm(0.1), Cm(30.0), Cm(1.8),
            "uc(y) = √[Σᵢ (∂f/∂xᵢ)² · u²(xᵢ)]",
            font_size=22, bold=True, color=COR_AZUL_ABAR, align=PP_ALIGN.CENTER)

y += Cm(2.5)

detalhes = [
    ("Componentes calculadas", "Incerteza tipo A (estatística, n=183 dias) e tipo B (calibração, resolução, temperatura)"),
    ("Grandezas de entrada", "Volume bruto, Poder Calorífico Superior (PCS), pressão, temperatura e fator de compressibilidade"),
    ("Resultado", "Incerteza expandida U = k · uc(y) para cada dia, com fator de abrangência k=2 (95% de confiança)"),
    ("Automação", "IA aplicou as equações da apostila em código Python para todos os 183 dias simultaneamente"),
]

for titulo, desc in detalhes:
    add_rect(slide, Cm(0.8), y, Cm(5.5), Cm(1.5), COR_AZUL_ABAR)
    add_textbox(slide, Cm(0.85), y + Cm(0.1), Cm(5.4), Cm(1.3),
                titulo, font_size=10, bold=True, color=COR_BRANCO,
                align=PP_ALIGN.CENTER, word_wrap=True)
    add_rect(slide, Cm(6.4), y, Cm(26.4), Cm(1.5), COR_CINZA_CLARO, COR_CINZA_MEDIO, Pt(1))
    add_textbox(slide, Cm(6.7), y + Cm(0.1), Cm(26.0), Cm(1.3),
                desc, font_size=11, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(1.75)

add_rect(slide, Cm(0.8), y + Cm(0.2), Cm(32.0), Cm(0.9), COR_VERDE_DESTAQUE)
add_textbox(slide, Cm(1.0), y + Cm(0.3), Cm(31.5), Cm(0.75),
            "O que levaria dias de planilha → 2 minutos de código gerado por IA",
            font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 14 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Balanço de Massa com IA
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Notebook 7: Balanço de Massa",
                                 subtitle_text="Verificação automática de consistência para 183 dias de operação")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

# Conceito
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(1.9), COR_CINZA_CLARO, COR_AZUL_ABAR, Pt(1))
add_textbox(slide, Cm(1.0), y + Cm(0.1), Cm(31.5), Cm(0.65),
            "CONCEITO DE BALANÇO DE MASSA:", font_size=12, bold=True, color=COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.8), Cm(31.5), Cm(0.95),
            "Volume entrada = Σ(volumes saída por cliente) ± incerteza de medição",
            font_size=14, bold=True, color=COR_AZUL_CLARO, word_wrap=True)

y += Cm(2.5)

# O que a IA fez
aspectos = [
    ("Leitura dos dados", "Leu volumes de entrada e todos os clientes (residencial, comercial, industrial, GNC) para cada um dos 183 dias"),
    ("Verificação automática", "Calculou o erro de balanço diário = [(Entrada - Σ Saídas) / Entrada] × 100%"),
    ("Alertas inteligentes", "Identificou automaticamente dias com desvio > 2% (limite aceitável segundo norma)"),
    ("Diagnóstico", "IA gerou texto explicativo para cada dia fora do limite — causa provável e recomendações"),
]

for titulo, desc in aspectos:
    add_rect(slide, Cm(0.8), y, Cm(5.8), Cm(1.6), COR_AZUL_ABAR)
    add_textbox(slide, Cm(0.9), y + Cm(0.1), Cm(5.6), Cm(1.4),
                titulo, font_size=10, bold=True, color=COR_BRANCO,
                align=PP_ALIGN.CENTER, word_wrap=True)
    add_rect(slide, Cm(6.7), y, Cm(26.1), Cm(1.6), COR_CINZA_CLARO, COR_CINZA_MEDIO, Pt(1))
    add_textbox(slide, Cm(7.0), y + Cm(0.1), Cm(25.7), Cm(1.4),
                desc, font_size=11, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(1.85)

# Conexão com o curso
add_rect(slide, Cm(0.8), y + Cm(0.2), Cm(32.0), Cm(1.1), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.3), Cm(31.5), Cm(0.9),
            "Conexão: mesmo verificador que o Prof. Barateiro apresentou nos Blocos 1-4 — agora automatizado para execução contínua",
            font_size=12, bold=True, color=COR_BRANCO, word_wrap=True)

print("  Slide 15 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Avaliação e Deployment
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Etapas 6-7: Avaliação e Deployment",
                                 subtitle_text="IBM Data Science Methodology • Passos 6-7 de 7")
add_logo(slide)
add_footer(slide)

# Pipeline visual
y = content_bottom + Cm(0.5)

etapas_pipeline = [
    ("Excel\n2,1 MB", COR_AZUL_ABAR),
    ("7 Notebooks\nJupyter", COR_AZUL_CLARO),
    ("28 chamadas\nGemini AI", RGBColor(0x4A, 0x86, 0xC8)),
    ("23 gráficos\ngerados", COR_VERDE_DESTAQUE),
    ("Relatório Word\n8,9 MB", COR_LARANJA),
]

bw = Cm(5.5)
bh = Cm(2.5)
xp = Cm(0.8)
yp = y

for i, (label, cor) in enumerate(etapas_pipeline):
    add_rect(slide, xp, yp, bw, bh, cor)
    add_textbox(slide, xp + Cm(0.1), yp + Cm(0.5),
                bw - Cm(0.2), bh - Cm(0.6),
                label, font_size=12, bold=True,
                color=COR_BRANCO, align=PP_ALIGN.CENTER, word_wrap=True)
    if i < len(etapas_pipeline) - 1:
        add_textbox(slide, xp + bw + Cm(0.05), yp + bh / 2 - Cm(0.3),
                    Cm(0.65), Cm(0.7), "→", font_size=22, bold=True,
                    color=COR_AZUL_ABAR, align=PP_ALIGN.CENTER)
    xp += bw + Cm(0.7)

# Tempo em destaque
y_tempo = yp + bh + Cm(0.6)
add_rect(slide, Cm(0.8), y_tempo, Cm(32.0), Cm(2.0), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y_tempo + Cm(0.1), Cm(31.5), Cm(0.75),
            "TEMPO TOTAL DE EXECUÇÃO:", font_size=12, bold=True, color=COR_CINZA_MEDIO)
add_textbox(slide, Cm(1.0), y_tempo + Cm(0.85), Cm(31.5), Cm(1.0),
            "1 minuto e 42 segundos  (execução paralela com 28 threads)",
            font_size=22, bold=True, color=COR_AMARELO_DEMO, align=PP_ALIGN.CENTER)

# Imagem download
y_img = y_tempo + Cm(2.6)
add_image_safe(slide, IMG_WEB6,
               Cm(0.8), y_img,
               Cm(18.0), Cm(8.0))

# Stats
stats = [
    ("8,9 MB", "Tamanho do\nrelatório final"),
    ("100+", "Páginas no\ndocumento Word"),
    ("23", "Gráficos\nincorporados"),
    ("3", "Diagramas\ngerados por IA"),
]
sx = Cm(20.0)
sy = y_img
sw = Cm(6.5)
sh = Cm(3.5)
gap_s = Cm(0.35)

for i, (num, desc) in enumerate(stats):
    xi = sx + (i % 2) * (sw + gap_s)
    yi = sy + (i // 2) * (sh + gap_s)
    add_rect(slide, xi, yi, sw, sh, COR_CINZA_CLARO, COR_AZUL_ABAR, Pt(1))
    add_textbox(slide, xi + Cm(0.2), yi + Cm(0.2), sw - Cm(0.4), Cm(1.3),
                num, font_size=26, bold=True, color=COR_AZUL_ABAR,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, xi + Cm(0.2), yi + Cm(1.6), sw - Cm(0.4), Cm(1.7),
                desc, font_size=10, color=COR_CINZA_TEXTO,
                align=PP_ALIGN.CENTER, word_wrap=True)

print("  Slide 16 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DEMO 2: Os 7 Notebooks ao Vivo
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, RGBColor(0xFF, 0xFF, 0xF0))
content_bottom = add_header_bar(slide, "DEMO AO VIVO — Os 7 Notebooks Jupyter",
                                 subtitle_text="Módulo 3 • Demo 2 de 4",
                                 bar_color=COR_LARANJA)
add_logo(slide)
add_footer(slide)
add_demo_banner(slide, content_bottom + Cm(0.3))

y = content_bottom + Cm(1.6)

notebooks = [
    ("01", "Leitura e Exploração do Excel", "Carrega todas as 14 abas, exibe estatísticas básicas e valida integridade dos dados"),
    ("02", "Análise de Volumes de Entrada", "Série temporal do volume diário no ponto de entrada do Distrito, tendências e sazonalidade"),
    ("03", "Análise de PCS", "Poder Calorífico Superior diário, variação por componente e comparação com limites normativos"),
    ("04", "Cálculo de Energia", "Energia por cliente (Volume × PCS), ranking de consumo e perfil mensal"),
    ("05", "Perfis de Consumo por Cliente", "Análise segmentada por categoria: residencial, comercial, industrial e GNC"),
    ("06", "Incertezas de Medição (GUM)", "Propagação de incerteza tipo A+B para cada grandeza, conforme apostila do Prof. Barateiro"),
    ("07", "Balanço de Massa", "Verificação Entrada = Σ Saídas, alertas de desvio e diagnóstico automático por IA"),
]

for num, titulo, desc in notebooks:
    num_color = [COR_AZUL_ABAR, COR_AZUL_CLARO, COR_VERDE_DESTAQUE, COR_LARANJA,
                 RGBColor(0x8B, 0x00, 0x8B), COR_AZUL_ABAR, COR_VERDE_DESTAQUE][int(num) - 1]
    add_rect(slide, Cm(0.8), y, Cm(1.2), Cm(1.0), num_color)
    add_textbox(slide, Cm(0.8), y, Cm(1.2), Cm(1.0),
                num, font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(2.3), y, Cm(9.5), Cm(0.5),
                titulo, font_size=12, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, Cm(2.3), y + Cm(0.5), Cm(30.5), Cm(0.5),
                desc, font_size=10, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(1.25)

# Nota final
add_rect(slide, Cm(0.8), y + Cm(0.2), Cm(32.0), Cm(0.9), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.3), Cm(31.5), Cm(0.75),
            "★  Nenhum notebook foi escrito manualmente — todos gerados e refinados com Claude Code",
            font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 17 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Arquitetura do Sistema
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Arquitetura: Do Excel ao Relatório em 1:42",
                                 subtitle_text="Visão completa do pipeline de auditoria automatizada")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

# Diagrama de fluxo (desenhado com formas)
# Linha 1: Excel → 7 Notebooks → 23 Gráficos
componentes_l1 = [
    ("Excel\n14 abas", COR_AZUL_ABAR, Cm(4.5)),
    ("7 Notebooks\nPython", COR_AZUL_CLARO, Cm(4.5)),
    ("23 Gráficos\nMatplotlib", COR_VERDE_DESTAQUE, Cm(4.5)),
]
xc = Cm(0.8)
yc = y
for label, cor, w in componentes_l1:
    add_rect(slide, xc, yc, w, Cm(1.8), cor)
    add_textbox(slide, xc + Cm(0.1), yc + Cm(0.3),
                w - Cm(0.2), Cm(1.2), label, font_size=11, bold=True,
                color=COR_BRANCO, align=PP_ALIGN.CENTER, word_wrap=True)
    xc += w + Cm(0.5)
    if label != componentes_l1[-1][0]:
        add_textbox(slide, xc - Cm(0.4), yc + Cm(0.5), Cm(0.45), Cm(0.8),
                    "→", font_size=18, bold=True, color=COR_AZUL_ABAR)

# Linha 2: Apostila → Gemini → 28 seções .md
yc2 = y + Cm(2.8)
componentes_l2 = [
    ("Apostila\nPDF", RGBColor(0x8B, 0x45, 0x13), Cm(4.5)),
    ("Gemini AI\nextrai metodologia", RGBColor(0x15, 0x65, 0xC0), Cm(5.5)),
    ("28 seções\n.md geradas", RGBColor(0x4A, 0x86, 0xC8), Cm(4.5)),
]
xc = Cm(0.8)
for label, cor, w in componentes_l2:
    add_rect(slide, xc, yc2, w, Cm(1.8), cor)
    add_textbox(slide, xc + Cm(0.1), yc2 + Cm(0.3),
                w - Cm(0.2), Cm(1.2), label, font_size=11, bold=True,
                color=COR_BRANCO, align=PP_ALIGN.CENTER, word_wrap=True)
    xc += w + Cm(0.5)
    if label != componentes_l2[-1][0]:
        add_textbox(slide, xc - Cm(0.4), yc2 + Cm(0.5), Cm(0.45), Cm(0.8),
                    "→", font_size=18, bold=True, color=COR_AZUL_ABAR)

# Seta de junção
yfinal = y + Cm(5.6)
add_textbox(slide, Cm(7.5), y + Cm(1.8), Cm(0.8), Cm(1.0),
            "↓", font_size=22, bold=True, color=COR_AZUL_ABAR)
add_textbox(slide, Cm(7.5), yc2 + Cm(1.8), Cm(0.8), Cm(1.0),
            "↓", font_size=22, bold=True, color=COR_AZUL_ABAR)

# Resultado final
add_rect(slide, Cm(4.5), yfinal, Cm(9.0), Cm(2.0), COR_LARANJA)
add_textbox(slide, Cm(4.6), yfinal + Cm(0.1), Cm(8.8), Cm(0.8),
            "DOCX Builder", font_size=14, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
add_textbox(slide, Cm(4.6), yfinal + Cm(0.9), Cm(8.8), Cm(0.9),
            "Relatório Word 8,9 MB\n100+ páginas", font_size=11, color=COR_BRANCO,
            align=PP_ALIGN.CENTER, word_wrap=True)

# Imagem do diagrama
add_image_safe(slide, IMG_PROCESSO,
               Cm(19.0), content_bottom + Cm(0.5),
               Cm(14.5), Cm(12.5))

print("  Slide 18 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — DEMO 3: Interface Web ao Vivo
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, RGBColor(0xFF, 0xFF, 0xF0))
content_bottom = add_header_bar(slide, "DEMO AO VIVO — Interface Web FastAPI",
                                 subtitle_text="Módulo 3 • Demo 3 de 4",
                                 bar_color=COR_LARANJA)
add_logo(slide)
add_footer(slide)
add_demo_banner(slide, content_bottom + Cm(0.3))

# Grid 2x3 de screenshots
web_imgs = [
    (IMG_WEB1, "01 Configuração"),
    (IMG_WEB2, "02 Pipeline"),
    (IMG_WEB3, "03 Gráficos"),
    (IMG_WEB4, "04 Diagramas"),
    (IMG_WEB5, "05 Textos LLM"),
    (IMG_WEB6, "06 Downloads"),
]

img_w = Cm(10.4)
img_h = Cm(6.5)
gap_x = Cm(0.5)
gap_y = Cm(0.4)
x_start = Cm(0.7)
y_start = content_bottom + Cm(1.5)

for i, (img_path, label) in enumerate(web_imgs):
    col = i % 3
    row = i // 3
    xi = x_start + col * (img_w + gap_x)
    yi = y_start + row * (img_h + Cm(0.5))

    add_image_safe(slide, img_path, xi, yi, img_w, img_h)
    add_rect(slide, xi, yi + img_h - Cm(0.6), img_w, Cm(0.6), RGBColor(0x00, 0x00, 0x00))
    add_textbox(slide, xi + Cm(0.1), yi + img_h - Cm(0.55), img_w - Cm(0.2), Cm(0.5),
                label, font_size=9, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

# Roteiro
add_rect(slide, Cm(0.7), SLIDE_H - Cm(2.0), Cm(32.2), Cm(0.8), COR_AZUL_ABAR)
add_textbox(slide, Cm(0.9), SLIDE_H - Cm(1.95), Cm(32.0), Cm(0.7),
            "Roteiro: Upload da planilha → progresso em tempo real → visualizar gráficos → baixar relatório",
            font_size=11, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 19 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — O Relatório Final
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "O Resultado: Relatório de Auditoria Técnica",
                                 subtitle_text="Output final do sistema — pronto para entrega ao regulador")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

# Stats em cards
stats = [
    ("8,9 MB", "Relatório gerado", COR_AZUL_ABAR),
    ("100+", "Páginas totais", COR_AZUL_CLARO),
    ("7", "Capítulos técnicos", COR_VERDE_DESTAQUE),
    ("23", "Gráficos gerados", COR_LARANJA),
    ("3", "Diagramas por IA", RGBColor(0x8B, 0x00, 0x8B)),
    ("1:42", "Tempo de geração", RGBColor(0xCC, 0x33, 0x00)),
]

card_w = Cm(5.0)
card_h = Cm(4.0)
x = Cm(0.8)
gap = Cm(0.35)

for i, (num, desc, cor) in enumerate(stats):
    add_rect(slide, x, y, card_w, card_h, COR_CINZA_CLARO, cor, Pt(2))
    add_rect(slide, x, y, card_w, Cm(0.4), cor)
    add_textbox(slide, x + Cm(0.1), y + Cm(0.5), card_w - Cm(0.2), Cm(1.8),
                num, font_size=30, bold=True, color=cor, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Cm(0.1), y + Cm(2.5), card_w - Cm(0.2), Cm(1.3),
                desc, font_size=10, color=COR_CINZA_TEXTO,
                align=PP_ALIGN.CENTER, word_wrap=True)
    x += card_w + gap

# Estrutura do relatório
y_struct = y + card_h + Cm(0.7)
add_rect(slide, Cm(0.8), y_struct, Cm(32.0), Cm(0.65), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y_struct + Cm(0.05), Cm(31.5), Cm(0.6),
            "Estrutura do relatório:", font_size=12, bold=True, color=COR_BRANCO)

y_struct += Cm(0.85)
estrutura = [
    "Cap. 1: Análise de Volumes de Entrada  |  Cap. 2: Poder Calorífico Superior (PCS)  |  Cap. 3: Energia Total",
    "Cap. 4: Perfis de Consumo por Cliente  |  Cap. 5: Incertezas de Medição (GUM)  |  Cap. 6: Balanço de Massa",
    "Cap. 7: Conclusões e Recomendações  |  Apêndice A: Metodologia  |  Apêndice B: Código Python dos 7 Notebooks",
]
for linha in estrutura:
    add_textbox(slide, Cm(0.8), y_struct, Cm(32.0), Cm(0.75),
                linha, font_size=10.5, color=COR_CINZA_TEXTO, word_wrap=True)
    y_struct += Cm(0.85)

add_rect(slide, Cm(0.8), y_struct + Cm(0.2), Cm(32.0), Cm(0.9), COR_VERDE_DESTAQUE)
add_textbox(slide, Cm(1.0), y_struct + Cm(0.3), Cm(31.5), Cm(0.75),
            "Mesmo conteúdo que levaria semanas de análise manual — gerado em 1 minuto e 42 segundos",
            font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 20 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Comparativo: Manual vs IA
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "O Mesmo Trabalho — Perspectivas Diferentes",
                                 subtitle_text="Análise complementar: conhecimento especializado + automação por IA")
add_logo(slide)
add_footer(slide)

# Tabela comparativa
col_widths = [Cm(7.0), Cm(12.0), Cm(12.5)]
col_labels = ["Aspecto", "Prof. Barateiro (Blocos 1-4)", "Este Sistema (Bloco 6)"]
rows = [
    ("Dados de entrada", "Planilha Excel — análise manual", "Planilha Excel — leitura automatizada"),
    ("Metodologia", "GUM, normas ABNT (domínio especializado)", "GUM extraída por Gemini da própria apostila"),
    ("Resultado", "Análise especializada e aprofundada", "Relatório de 100+ páginas em 1:42 min"),
    ("Programação", "Especialistas em Python", "Zero linha de Python manual"),
    ("Interpretação", "Especialista humano — contexto e julgamento", "IA gera texto técnico — revisão humana necessária"),
    ("Escalabilidade", "Limitada — requer tempo especializado", "Automático — re-executa em qualquer base"),
]

table_x = Cm(0.8)
table_y = content_bottom + Cm(0.5)
row_h = Cm(1.7)

x = table_x
for j, (label, w) in enumerate(zip(col_labels, col_widths)):
    cor_hdr = COR_AZUL_ABAR if j == 0 else ([COR_AZUL_CLARO, COR_VERDE_DESTAQUE][j - 1])
    add_rect(slide, x, table_y, w, row_h, cor_hdr, COR_BRANCO, Pt(1))
    add_textbox(slide, x + Cm(0.2), table_y + Cm(0.2), w - Cm(0.4), row_h - Cm(0.4),
                label, font_size=11, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER,
                word_wrap=True)
    x += w

for i, row in enumerate(rows):
    y_row = table_y + (i + 1) * row_h
    x = table_x
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        bg = COR_CINZA_CLARO if i % 2 == 0 else COR_BRANCO
        add_rect(slide, x, y_row, w, row_h, bg, COR_CINZA_MEDIO, Pt(1))
        bold = (j == 0)
        text_color = COR_AZUL_ABAR if j == 0 else COR_CINZA_TEXTO
        add_textbox(slide, x + Cm(0.2), y_row + Cm(0.15), w - Cm(0.4), row_h - Cm(0.3),
                    cell, font_size=10, bold=bold, color=text_color,
                    align=PP_ALIGN.LEFT, word_wrap=True)
        x += w

# Mensagem final
bottom_y = table_y + (len(rows) + 1) * row_h + Cm(0.4)
add_rect(slide, Cm(0.8), bottom_y, Cm(31.5), Cm(1.0), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), bottom_y + Cm(0.1), Cm(31.2), Cm(0.85),
            "Complementar, não concorrente — a IA amplifica o conhecimento técnico do especialista humano",
            font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 21 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Dashboard e Publicação Web
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Indo Além: Dashboard e Publicação",
                                 subtitle_text="Com IA, criar um dashboard não requer mais saber programar")
add_logo(slide)
add_footer(slide)

# Imagem lateral
add_image_safe(slide, IMG_WEB3,
               Cm(19.5), content_bottom + Cm(0.5),
               Cm(13.8), Cm(9.5))

y = content_bottom + Cm(0.5)

opcoes = [
    ("Streamlit", "streamlit.io", "App Python simples, deploy gratuito em minutos. Ideal para protótipos e projetos internos."),
    ("FastAPI + HTML", "fastapi.tiangolo.com", "Como o sistema demonstrado hoje. Alta performance, interface customizável."),
    ("Google Looker Studio", "lookerstudio.google.com", "Totalmente sem código. Conecta diretamente ao Excel ou Google Sheets. Gratuito."),
    ("Power BI / Tableau", "powerbi.com / tableau.com", "Para ambientes corporativos. Integração com sistemas legados e ERP."),
]

for plataforma, url, desc in opcoes:
    add_rect(slide, Cm(0.8), y, Cm(3.5), Cm(1.9), COR_AZUL_ABAR)
    add_textbox(slide, Cm(0.85), y + Cm(0.2), Cm(3.4), Cm(0.9),
                plataforma, font_size=11, bold=True, color=COR_BRANCO,
                align=PP_ALIGN.CENTER, word_wrap=True)
    add_textbox(slide, Cm(0.85), y + Cm(1.1), Cm(3.4), Cm(0.65),
                url, font_size=8, color=COR_AMARELO_DEMO, align=PP_ALIGN.CENTER)
    add_rect(slide, Cm(4.4), y, Cm(14.5), Cm(1.9), COR_CINZA_CLARO, COR_CINZA_MEDIO, Pt(1))
    add_textbox(slide, Cm(4.6), y + Cm(0.2), Cm(14.2), Cm(1.5),
                desc, font_size=10.5, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(2.2)

add_rect(slide, Cm(0.8), y + Cm(0.2), Cm(18.0), Cm(1.0), COR_VERDE_DESTAQUE)
add_textbox(slide, Cm(1.0), y + Cm(0.3), Cm(17.8), Cm(0.75),
            "Com IA: descreva o dashboard → IA escreve o código",
            font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 22 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — DEMO 4: NotebookLM para os Alunos
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, RGBColor(0xFF, 0xFF, 0xF0))
content_bottom = add_header_bar(slide, "DEMO AO VIVO — NotebookLM: Leve a Apostila para Conversar",
                                 subtitle_text="Módulo 4 • Demo 4 de 4",
                                 bar_color=COR_LARANJA)
add_logo(slide)
add_footer(slide)
add_demo_banner(slide, content_bottom + Cm(0.3))

y = content_bottom + Cm(1.6)

# Instruções
add_textbox(slide, Cm(0.8), y, Cm(32.0), Cm(0.7),
            "INSTRUCOES PARA OS ALUNOS:", font_size=13, bold=True, color=COR_AZUL_ABAR)
y += Cm(0.8)

instrucoes = [
    ("1", "Acesse notebooklm.google.com — faça login com conta Google (gratuito)"),
    ("2", "Clique em 'Criar novo notebook'"),
    ("3", "Clique em 'Adicionar fontes' → faça upload da apostila do Prof. Barateiro"),
    ("4", "Faça perguntas técnicas em linguagem natural no chat"),
]

for num, texto in instrucoes:
    add_rect(slide, Cm(0.8), y, Cm(1.2), Cm(1.0), COR_AZUL_ABAR)
    add_textbox(slide, Cm(0.8), y, Cm(1.2), Cm(1.0),
                num, font_size=14, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(2.3), y, Cm(30.5), Cm(1.0),
                texto, font_size=12, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(1.3)

y += Cm(0.2)
add_textbox(slide, Cm(0.8), y, Cm(32.0), Cm(0.7),
            "EXEMPLOS DE PERGUNTAS TECNICAS:", font_size=12, bold=True, color=COR_AZUL_ABAR)
y += Cm(0.8)

perguntas = [
    "\"Qual a diferença entre incerteza tipo A e tipo B segundo a apostila?\"",
    "\"Como calcular o PCS do gás natural a partir da composição cromatográfica?\"",
    "\"Quais são os requisitos de medição fiscal em dutos offshore segundo a norma?\"",
    "\"Explique o método GUM de propagação de incertezas em termos simples\"",
]

for pergunta in perguntas:
    add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(0.95), RGBColor(0xE8, 0xF4, 0xFD), COR_AZUL_CLARO, Pt(1))
    add_textbox(slide, Cm(1.2), y + Cm(0.1), Cm(31.5), Cm(0.8),
                pergunta, font_size=11, color=COR_AZUL_ABAR, word_wrap=True,
                font_name=FONTE_CORPO)
    y += Cm(1.1)

print("  Slide 23 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Materiais e Próximos Passos
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_BRANCO)
content_bottom = add_header_bar(slide, "Materiais para Continuar Aprendendo",
                                 subtitle_text="Recursos gratuitos para levar para casa")
add_logo(slide)
add_footer(slide)

y = content_bottom + Cm(0.5)

# Seção 1
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(0.65), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.05), Cm(31.5), Cm(0.6),
            "Ferramentas para começar HOJE (todas gratuitas):",
            font_size=12, bold=True, color=COR_BRANCO)
y += Cm(0.85)

ferramentas_gratis = [
    ("chatgpt.com", "Análise de planilhas, gráficos, resumos — sem instalar nada"),
    ("notebooklm.google.com", "Conversas com documentos técnicos — carregar apostila"),
    ("github.com/copilot", "Assistência de código — freemium, 2000 sugestões/mês grátis"),
]

for url, desc in ferramentas_gratis:
    add_rect(slide, Cm(0.8), y, Cm(6.5), Cm(0.95), COR_AZUL_CLARO)
    add_textbox(slide, Cm(0.85), y + Cm(0.1), Cm(6.4), Cm(0.8),
                url, font_size=10, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(7.6), y, Cm(25.0), Cm(0.95),
                "→  " + desc, font_size=11, color=COR_CINZA_TEXTO)
    y += Cm(1.1)

y += Cm(0.3)

# Seção 2
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(0.65), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.05), Cm(31.5), Cm(0.6),
            "Artigos recomendados (acesso gratuito):",
            font_size=12, bold=True, color=COR_BRANCO)
y += Cm(0.85)

artigos = [
    ("arxiv.org/abs/2303.12712", '"Sparks of AGI" — Microsoft Research — 154 páginas sobre capacidades emergentes do GPT-4'),
    ("arxiv.org/abs/2305.13657", '"ChatGPT as Your Personal Data Scientist" — análise de dados com LLMs'),
    ("naruc.org", '"AI for Gas Utilities" — Guia NARUC para reguladores de gás e energia'),
]

for url, desc in artigos:
    add_rect(slide, Cm(0.8), y, Cm(6.5), Cm(0.95), COR_VERDE_DESTAQUE)
    add_textbox(slide, Cm(0.85), y + Cm(0.1), Cm(6.4), Cm(0.8),
                url, font_size=8, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER, word_wrap=True)
    add_textbox(slide, Cm(7.6), y, Cm(25.0), Cm(0.95),
                desc, font_size=10.5, color=COR_CINZA_TEXTO, word_wrap=True)
    y += Cm(1.1)

y += Cm(0.3)

# Seção 3
add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(0.65), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.0), y + Cm(0.05), Cm(31.5), Cm(0.6),
            "Este projeto:",
            font_size=12, bold=True, color=COR_BRANCO)
y += Cm(0.85)

add_rect(slide, Cm(0.8), y, Cm(32.0), Cm(1.5), COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(1))
add_textbox(slide, Cm(1.0), y + Cm(0.1), Cm(31.5), Cm(1.3),
            "• 7 Notebooks Jupyter disponibilizados para os alunos\n"
            "• Artigos organizados por tema na pasta 'artigos/'\n"
            "• Contato: raul.araujo@agenersa.rj.gov.br",
            font_size=11, color=COR_CINZA_TEXTO, word_wrap=True)

print("  Slide 24 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — AVALIAÇÃO
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
setup_slide_background(slide, COR_AZUL_ABAR)

# Faixa decorativa
add_rect(slide, Cm(0), Cm(0), Cm(0.6), SLIDE_H, COR_AZUL_CLARO)

# Painel branco
painel_y = Cm(1.8)
painel_h = SLIDE_H - painel_y - Cm(1.8)
add_rect(slide, Cm(0.6), painel_y, SLIDE_W - Cm(0.6) - Cm(1.5), painel_h, COR_BRANCO)

# Logo
if os.path.exists(IMG_LOGO):
    slide.shapes.add_picture(IMG_LOGO, Cm(1.2), Cm(0.2), Cm(6.0), Cm(1.8))

# Título
add_textbox(slide, Cm(1.5), painel_y + Cm(0.5), Cm(30.0), Cm(2.0),
            "Avaliação — Bloco 6",
            font_size=32, bold=True, color=COR_AZUL_ABAR, align=PP_ALIGN.CENTER)

add_textbox(slide, Cm(1.5), painel_y + Cm(2.5), Cm(30.0), Cm(1.0),
            "10 questões de múltipla escolha",
            font_size=18, bold=False, color=COR_AZUL_CLARO, align=PP_ALIGN.CENTER)

# Linha divisória
add_rect(slide, Cm(4.0), painel_y + Cm(3.7), Cm(24.5), Cm(0.06), COR_AZUL_CLARO)

# Instruções
add_textbox(slide, Cm(1.5), painel_y + Cm(4.0), Cm(30.0), Cm(1.5),
            "Documento entregue em papel ou via link digital\n"
            "Responda individualmente — tempo: 15 minutos",
            font_size=14, color=COR_CINZA_TEXTO, align=PP_ALIGN.CENTER, word_wrap=True)

# Placeholder QR Code
add_rect(slide, Cm(13.0), painel_y + Cm(5.8), Cm(4.5), Cm(4.5), COR_CINZA_CLARO, COR_AZUL_ABAR, Pt(2))
add_textbox(slide, Cm(13.0), painel_y + Cm(7.0), Cm(4.5), Cm(1.5),
            "[QR Code\nAvaliação]", font_size=10, color=COR_AZUL_ABAR, align=PP_ALIGN.CENTER, word_wrap=True)

# Mensagem final
add_textbox(slide, Cm(1.5), painel_y + Cm(10.8), Cm(30.0), Cm(1.5),
            "Obrigado! Dúvidas e comentários são bem-vindos.",
            font_size=18, bold=True, color=COR_AZUL_ABAR, align=PP_ALIGN.CENTER)

# Rodapé azul
add_rect(slide, Cm(0.6), SLIDE_H - Cm(1.8), SLIDE_W - Cm(2.1), Cm(1.8), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.2), SLIDE_H - Cm(1.6), Cm(30.0), Cm(0.9),
            "raul.araujo@agenersa.rj.gov.br  |  Curso ABAR — Fevereiro 2026",
            font_size=12, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 25 OK")


# ─── Salvar ───────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PPTX)
size = os.path.getsize(OUTPUT_PPTX)
print(f"\n{'='*60}")
print(f"SUCESSO! Apresentação criada:")
print(f"  Arquivo: {OUTPUT_PPTX}")
print(f"  Tamanho: {size:,} bytes ({size / 1024 / 1024:.1f} MB)")
print(f"  Slides: 25")
print(f"{'='*60}")
