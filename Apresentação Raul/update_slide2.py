"""
Script para atualizar o Slide 2 do arquivo PPTX existente com novos benchmarks.
Mantém o estilo visual (azul #003886, fonte Calibri, logo ABAR, rodapé) dos outros slides.

Uso: python update_slide2.py
"""

import os
import sys
import copy
from lxml import etree

try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Caminhos ─────────────────────────────────────────────────────────────────
BASE = r'C:\Users\raula\OneDrive\Documentos\Codigos e sistemas pessoais\AGENERSA\Cursos ABAR de Dados\Apresentação Raul'
SCREENS = os.path.join(BASE, 'screenshots')

PPTX_FILE = os.path.join(BASE, 'Bloco6_AI_Data_Science_ABAR.pptx')
IMG_LOGO   = os.path.join(BASE, 'template_Imagem 10.jpg')
IMG_TRACKINGAI = os.path.join(SCREENS, 'benchmark_trackingai_home.png')

# ─── Cores ────────────────────────────────────────────────────────────────────
COR_AZUL_ABAR   = RGBColor(0x00, 0x38, 0x86)   # #003886
COR_AZUL_CLARO  = RGBColor(0x00, 0x81, 0xBF)   # #0081BF
COR_BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
COR_CINZA_TEXTO = RGBColor(0x44, 0x44, 0x55)
COR_VERDE       = RGBColor(0x00, 0x8A, 0x4B)
COR_LARANJA     = RGBColor(0xE8, 0x6D, 0x00)
COR_CINZA_CLARO = RGBColor(0xF0, 0xF4, 0xF8)
COR_CINZA_MEDIO = RGBColor(0xD0, 0xD8, 0xE0)

# ─── Dimensões (widescreen 16:9) ──────────────────────────────────────────────
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

FONTE_TITULO = "Calibri"
FONTE_CORPO  = "Calibri"

# ─── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=COR_CINZA_TEXTO, align=PP_ALIGN.LEFT,
                font_name=FONTE_CORPO, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = word_wrap
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return txBox


def add_image_safe(slide, img_path, left, top, width, height):
    if os.path.exists(img_path):
        return slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        shape = add_rect(slide, left, top, width, height, COR_CINZA_MEDIO, COR_AZUL_CLARO)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "[Imagem: %s]" % os.path.basename(img_path)
        run.font.size = Pt(9)
        run.font.color.rgb = COR_CINZA_TEXTO
        return shape


def remove_slide_shapes(slide):
    """Remove todas as shapes de conteúdo do slide (mantém apenas elementos do layout/master)."""
    sp_tree = slide.shapes._spTree
    # Coletar todos os elementos de shape para remover
    to_remove = []
    for child in sp_tree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        # Remover: sp (textbox/shape), pic (picture), grpSp (group)
        # Manter: spTree metadata (nvGrpSpPr, grpSpPr, etc.)
        if tag in ('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'):
            to_remove.append(child)
    for elem in to_remove:
        sp_tree.remove(elem)


# ─── Abrir PPTX ───────────────────────────────────────────────────────────────
print("Abrindo: %s" % PPTX_FILE)
prs = Presentation(PPTX_FILE)

print("Dimensoes: %.2f x %.2f inches" % (prs.slide_width.inches, prs.slide_height.inches))
print("Total de slides: %d" % len(prs.slides))

# ─── Acessar Slide 2 (índice 1) ───────────────────────────────────────────────
slide = prs.slides[1]
print("Slide 2 atual: %d shapes" % len(slide.shapes))

# Remover todas as shapes existentes
remove_slide_shapes(slide)
print("Shapes removidas. Shapes restantes: %d" % len(slide.shapes))

# ─── Fundo branco ─────────────────────────────────────────────────────────────
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = COR_BRANCO

# ─── Header bar azul ──────────────────────────────────────────────────────────
BAR_H = Cm(3.0)
add_rect(slide, Cm(0), Cm(0), SLIDE_W, BAR_H, COR_AZUL_ABAR)

# Título principal no header
title_box = slide.shapes.add_textbox(Cm(0.8), Cm(0.2), SLIDE_W - Cm(6.5), BAR_H - Cm(0.5))
title_box.text_frame.word_wrap = True
p = title_box.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "A IA Evoluiu Mais em 2 Anos do que em 20"
run.font.size = Pt(24)
run.font.bold = True
run.font.name = FONTE_TITULO
run.font.color.rgb = COR_BRANCO

# Subtítulo no header
sub_box = slide.shapes.add_textbox(Cm(0.8), BAR_H - Cm(0.95), SLIDE_W - Cm(6.5), Cm(0.85))
sub_box.text_frame.word_wrap = True
p_sub = sub_box.text_frame.paragraphs[0]
p_sub.alignment = PP_ALIGN.LEFT
run_sub = p_sub.add_run()
run_sub.text = "Benchmarks que mostram uma ruptura tecnológica sem precedentes"
run_sub.font.size = Pt(13)
run_sub.font.bold = False
run_sub.font.name = FONTE_CORPO
run_sub.font.color.rgb = COR_AZUL_CLARO

# ─── Logo ABAR no canto superior direito ──────────────────────────────────────
logo_w = Cm(5.0)
logo_h = Cm(1.5)
logo_left = SLIDE_W - logo_w - Cm(0.3)
logo_top = Cm(0.5)
if os.path.exists(IMG_LOGO):
    slide.shapes.add_picture(IMG_LOGO, logo_left, logo_top, logo_w, logo_h)

# ─── Área de conteúdo ─────────────────────────────────────────────────────────
CONTENT_TOP = BAR_H + Cm(0.35)

# Layout: 60% esquerda (bullets) | 40% direita (imagem)
LEFT_W  = SLIDE_W * 0.60   # ~20.3 cm
RIGHT_W = SLIDE_W * 0.40   # ~13.5 cm
RIGHT_X = LEFT_W + Cm(0.3)

FOOTER_H = Cm(0.70)
CONTENT_H = SLIDE_H - CONTENT_TOP - FOOTER_H - Cm(0.2)

# ─── COLUNA ESQUERDA — Bullets ────────────────────────────────────────────────

# Definir bullets com badge e texto
bullets = [
    ("Nov/2024", COR_AZUL_ABAR,
     "GPT-4o marcou apenas 3% no Humanity's Last Exam (HLE)"),
    ("Jan/2026", COR_VERDE,
     "Claude Opus 4.6 marca 40% no mesmo teste — em apenas 14 meses"),
    ("2024", COR_LARANJA,
     "Modelos IA com QI ~65 (abaixo da média humana de 100)"),
    ("2026", COR_VERDE,
     "Claude Opus 4.6 com QI ~130 — faixa superdotado (topo 2% humano)"),
    ("Jul/2025", COR_AZUL_ABAR,
     "Gemini Deep Think → Medalha de Ouro na Olimpíada Internacional de Matemática (IMO)"),
    ("Hoje", COR_AZUL_CLARO,
     "Resolvem 76,8% de bugs reais de software (SWE-bench)"),
    ("FrontierMath", COR_AZUL_ABAR,
     "Matemática nível PhD: de 0% → 40% em apenas 12 meses"),
    ("Reflexão", COR_VERDE,
     "\"Se chegou aqui... o que faz com dados de distribuição de gás?\""),
]

BADGE_W = Cm(3.2)
BADGE_H = Cm(0.68)
ROW_GAP = Cm(0.25)

# Calcular posição Y distribuindo os bullets na área disponível
total_rows = len(bullets)
row_h = BADGE_H + ROW_GAP
total_h = total_rows * row_h - ROW_GAP

# Centralizar verticalmente na coluna esquerda
y_start = CONTENT_TOP + (CONTENT_H - total_h) / 2
if y_start < CONTENT_TOP:
    y_start = CONTENT_TOP

y = y_start
for badge_text, badge_color, texto in bullets:
    # Badge colorido (retângulo arredondado via shape normal)
    badge = add_rect(slide, Cm(0.5), y, BADGE_W, BADGE_H, badge_color)
    add_textbox(slide, Cm(0.5), y + Cm(0.04),
                BADGE_W, BADGE_H - Cm(0.08),
                badge_text, font_size=9, bold=True,
                color=COR_BRANCO, align=PP_ALIGN.CENTER)

    # Texto do bullet
    txt_x = Cm(0.5) + BADGE_W + Cm(0.25)
    txt_w = LEFT_W - BADGE_W - Cm(1.2)
    is_reflexao = (badge_text == "Reflexão")
    txt_box = slide.shapes.add_textbox(txt_x, y - Cm(0.04), txt_w, BADGE_H + Cm(0.2))
    txt_box.text_frame.word_wrap = True
    p = txt_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(11.5)
    run.font.name = FONTE_CORPO
    run.font.bold = is_reflexao
    run.font.italic = is_reflexao
    run.font.color.rgb = COR_VERDE if is_reflexao else COR_CINZA_TEXTO

    y += row_h

# ─── COLUNA DIREITA — Imagem ──────────────────────────────────────────────────
IMG_TOP = CONTENT_TOP + Cm(0.2)
IMG_H   = CONTENT_H - Cm(0.4)
IMG_W   = SLIDE_W - RIGHT_X - Cm(0.4)

# Borda azul sutil ao redor da imagem
add_rect(slide, RIGHT_X - Cm(0.06), IMG_TOP - Cm(0.06),
         IMG_W + Cm(0.12), IMG_H + Cm(0.12),
         COR_AZUL_CLARO)

add_image_safe(slide, IMG_TRACKINGAI,
               RIGHT_X, IMG_TOP, IMG_W, IMG_H)

# ─── Rodapé ───────────────────────────────────────────────────────────────────
FOOTER_TOP = SLIDE_H - FOOTER_H
add_rect(slide, Cm(0), FOOTER_TOP, SLIDE_W, FOOTER_H, COR_AZUL_ABAR)
add_textbox(
    slide,
    Cm(0.5), FOOTER_TOP + Cm(0.05),
    SLIDE_W - Cm(1.0), FOOTER_H - Cm(0.1),
    "Fontes: TrackingAI.org | agi.safe.ai | epoch.ai/frontiermath   |   Curso ABAR — Medições Inteligentes e Gestão Integrada | Fevereiro 2026",
    font_size=8,
    color=COR_BRANCO,
    align=PP_ALIGN.CENTER
)

# ─── Salvar ───────────────────────────────────────────────────────────────────
print("Salvando em: %s" % PPTX_FILE)
prs.save(PPTX_FILE)
print("Slide 2 atualizado com sucesso!")
print("  Titulo: 'A IA Evoluiu Mais em 2 Anos do que em 20'")
print("  Bullets: %d itens com benchmarks" % len(bullets))
print("  Imagem: benchmark_trackingai_home.png (curva de QI)")
print("  Rodape com fontes: TrackingAI.org | agi.safe.ai | epoch.ai/frontiermath")
