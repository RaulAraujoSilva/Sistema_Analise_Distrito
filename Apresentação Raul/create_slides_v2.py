"""
Script v2 para criar a apresentação PPTX do Bloco 6 — IA aplicada à Ciência de Dados
Curso ABAR de Medições Inteligentes e Gestão Integrada | Fevereiro 2026

Melhorias v2:
- 28 slides (vs 25 no v1)
- Correção de alinhamento: posições calculadas, nunca hardcoded
- Screenshots reais de ferramentas
- Novos slides: 4 Modelos, Capacidades, IDEs, CLIs, Agentes, Claude+Excel, 2 Dashboards
- Removidos slides redundantes GUM/Balanço standalone

Uso: python create_slides_v2.py
"""

import os
import sys

# ─── Verificar/instalar python-pptx ───────────────────────────────────────────
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Caminhos ─────────────────────────────────────────────────────────────────
BASE = r'C:\Users\raula\OneDrive\Documentos\Codigos e sistemas pessoais\AGENERSA\Cursos ABAR de Dados\Apresentação Raul'
OUTPUTS = r'C:\Users\raula\OneDrive\Documentos\Codigos e sistemas pessoais\AGENERSA\Cursos ABAR de Dados\outputs'
SCREENS_PRES = os.path.join(BASE, 'screenshots')
SCREENS_TOOLS = os.path.join(BASE, 'screenshots', 'tools')
SCREENS_WEB = os.path.join(OUTPUTS, 'screenshots')
DIAGRAMAS = os.path.join(OUTPUTS, 'diagramas')

OUTPUT_PPTX = os.path.join(BASE, 'Bloco6_AI_Data_Science_ABAR_v2.pptx')

# ─── Funções auxiliares para caminhos de imagem ───────────────────────────────
def img(filename, folder=SCREENS_PRES):
    return os.path.join(folder, filename)

def img_tool(filename):
    return os.path.join(SCREENS_TOOLS, filename)

def img_web(filename):
    return os.path.join(SCREENS_WEB, filename)

def img_diag(filename):
    return os.path.join(DIAGRAMAS, filename)

# ─── Imagens ──────────────────────────────────────────────────────────────────
IMG_LOGO   = os.path.join(BASE, 'template_Imagem 10.jpg')
IMG_BG     = os.path.join(BASE, 'template_bg.png')
IMG_IBM    = os.path.join(BASE, 'WhatsApp Image 2026-02-12 at 10.32.47.jpeg')

# Benchmarks (já existentes)
IMG_IMO        = img('benchmark_imo_2025.png')
IMG_SWEBENCH   = img('benchmark_swebench.png')
IMG_EPOCH      = img('benchmark_epoch_ai.png')
IMG_CHATBOT    = img('benchmark_chatbot_arena.png')
IMG_ARC        = img('benchmark_arc_agi.png')
IMG_STANFORD   = img('benchmark_ai_timeline_stanford.png')
IMG_FRONTIER   = img('benchmark_frontiermath_leaderboard.png')
IMG_HLE        = img('benchmark_humanitys_last_exam_chart.png')
IMG_ARTIFICIAL = img('benchmark_artificialanalysis_intelligence.png')

# Ferramentas AI — capturados pelos subagentes
IMG_CHATGPT    = img_tool('chatgpt_interface.png')
IMG_GEMINI     = img_tool('gemini_interface.png')
IMG_CLAUDE     = img_tool('claude_interface.png')
IMG_GROK       = img_tool('grok_interface.png')
IMG_VSCODE     = img_tool('vscode_copilot.png')
IMG_CURSOR     = img_tool('cursor_ai.png')
IMG_WINDSURF   = img_tool('windsurf.png')
IMG_CODEX      = img_tool('codex_cli.png')
IMG_MANUS      = img_tool('manus_platform.png')
IMG_GENSPARK   = img_tool('genspark_platform.png')
IMG_COWORK     = img_tool('claude_cowork.png')

# Dados e dashboards
IMG_KAGGLE     = img('kaggle_housing.png')
IMG_GOVBR      = img('govbr_dados.png')
IMG_OBS_GAS    = img('obs_gas_dashboard.png')
IMG_ANP        = img('anp_dashboard.png')
IMG_STREAMLIT  = img('streamlit_example.png')

# Web screenshots do sistema
IMG_WEB1 = img_web('01_web_config.png')
IMG_WEB2 = img_web('02_web_pipeline.png')
IMG_WEB3 = img_web('03_web_graficos.png')
IMG_WEB4 = img_web('04_web_diagramas.png')
IMG_WEB5 = img_web('05_web_textos.png')
IMG_WEB6 = img_web('06_web_downloads.png')

# Diagramas
IMG_ESTRUTURA = img_diag('estrutura_distrito.png')
IMG_FLUXO     = img_diag('fluxo_auditoria.png')
IMG_PROCESSO  = img_diag('processo_analise.png')

# ─── Cores ────────────────────────────────────────────────────────────────────
COR_AZUL_ABAR      = RGBColor(0x00, 0x38, 0x86)
COR_AZUL_CLARO     = RGBColor(0x00, 0x81, 0xBF)
COR_BRANCO         = RGBColor(0xFF, 0xFF, 0xFF)
COR_CINZA_CLARO    = RGBColor(0xF0, 0xF4, 0xF8)
COR_CINZA_MEDIO    = RGBColor(0xD0, 0xD8, 0xE0)
COR_CINZA_TEXTO    = RGBColor(0x44, 0x44, 0x55)
COR_VERDE_DESTAQUE = RGBColor(0x00, 0x8A, 0x4B)
COR_LARANJA        = RGBColor(0xE8, 0x6D, 0x00)
COR_AMARELO_DEMO   = RGBColor(0xFF, 0xD7, 0x00)
COR_ROXO           = RGBColor(0x8B, 0x00, 0x8B)
COR_VERMELHO       = RGBColor(0xC0, 0x20, 0x20)

# ─── Dimensões (widescreen 16:9) ──────────────────────────────────────────────
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

# ─── Layout global (REGRA: nunca hardcode — calcular a partir destas constantes)
MARGIN      = Cm(0.5)
HEADER_H    = Cm(2.8)    # altura da barra de header (incluindo subtítulo)
FOOTER_H    = Cm(0.7)
CONTENT_TOP = HEADER_H + Cm(0.3)     # topo da área de conteúdo
CONTENT_H   = SLIDE_H - CONTENT_TOP - FOOTER_H - MARGIN  # altura útil ≈ 14.75cm
CONTENT_W   = SLIDE_W - MARGIN * 2   # largura útil ≈ 32.87cm
CONTENT_L   = MARGIN                  # margem esquerda

# ─── Fontes ───────────────────────────────────────────────────────────────────
FONTE_TITULO = "Calibri"
FONTE_CORPO  = "Calibri"

# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def add_image_safe(slide, img_path, left, top, width, height, label=None):
    """Adiciona imagem se existir, senão retângulo placeholder."""
    if img_path and os.path.exists(img_path):
        try:
            return slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception as e:
            print(f"    [WARN] Erro ao inserir {os.path.basename(img_path)}: {e}")
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COR_CINZA_MEDIO
    shape.line.color.rgb = COR_AZUL_CLARO
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    display = label or (os.path.basename(img_path) if img_path else "imagem")
    run.text = f"[{display}]"
    run.font.size = Pt(9)
    run.font.color.rgb = COR_CINZA_TEXTO
    return shape


def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=14, bold=False, color=COR_CINZA_TEXTO,
                align=PP_ALIGN.LEFT, font_name=FONTE_CORPO, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = word_wrap
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return txBox


def add_paragraph(text_frame, text, font_size=12, bold=False,
                  color=COR_CINZA_TEXTO, align=PP_ALIGN.LEFT,
                  font_name=FONTE_CORPO, space_before=0, indent_level=0):
    p = text_frame.add_paragraph()
    p.alignment = align
    p.level = indent_level
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return p


def setup_slide_background(slide, color=COR_BRANCO):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text, subtitle_text=None):
    """Adiciona barra de cabeçalho com altura fixa HEADER_H."""
    add_rect(slide, Cm(0), Cm(0), SLIDE_W, HEADER_H, COR_AZUL_ABAR)
    title_h = Cm(1.5) if subtitle_text else HEADER_H
    title_top = (HEADER_H - title_h) / 2 if not subtitle_text else Cm(0.15)
    txBox = slide.shapes.add_textbox(Cm(0.7), title_top, SLIDE_W - Cm(6.5), title_h)
    txBox.text_frame.word_wrap = True
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(22) if subtitle_text else Pt(24)
    run.font.bold = True
    run.font.name = FONTE_TITULO
    run.font.color.rgb = COR_BRANCO
    if subtitle_text:
        sub = slide.shapes.add_textbox(Cm(0.7), Cm(1.7), SLIDE_W - Cm(7), Cm(0.9))
        sub.text_frame.word_wrap = True
        p2 = sub.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(12)
        run2.font.name = FONTE_CORPO
        run2.font.color.rgb = COR_AZUL_CLARO


def add_footer(slide, text="Curso ABAR — Medições Inteligentes e Gestão Integrada  |  Fevereiro 2026"):
    footer_top = SLIDE_H - FOOTER_H
    add_rect(slide, Cm(0), footer_top, SLIDE_W, FOOTER_H, COR_AZUL_ABAR)
    add_textbox(slide, Cm(0.5), footer_top + Cm(0.05),
                SLIDE_W - Cm(1), FOOTER_H - Cm(0.1),
                text, font_size=8, color=COR_BRANCO, align=PP_ALIGN.CENTER)


def add_logo(slide):
    logo_w, logo_h = Cm(5.0), Cm(1.5)
    logo_left = SLIDE_W - logo_w - Cm(0.3)
    logo_top = Cm(0.3)
    if os.path.exists(IMG_LOGO):
        slide.shapes.add_picture(IMG_LOGO, logo_left, logo_top, logo_w, logo_h)


def add_demo_banner(slide, top_offset=None):
    if top_offset is None:
        top_offset = CONTENT_TOP
    banner_w = Cm(10)
    add_rect(slide, CONTENT_L, top_offset, banner_w, Cm(0.9),
             COR_AMARELO_DEMO, COR_LARANJA, Pt(2))
    add_textbox(slide, CONTENT_L, top_offset, banner_w, Cm(0.9),
                "★  DEMO AO VIVO", font_size=13, bold=True,
                color=RGBColor(0x80, 0x30, 0x00), align=PP_ALIGN.CENTER)


def add_card_tool(slide, left, top, width, height, nome, url, descricao, cor, img_path=None):
    """Card para slide de ferramentas com screenshot opcional."""
    # Fundo do card
    add_rect(slide, left, top, width, height, COR_CINZA_CLARO, cor, Pt(2))
    # Header do card
    header_h = Cm(1.1)
    add_rect(slide, left, top, width, header_h, cor)
    add_textbox(slide, left + Cm(0.2), top + Cm(0.05),
                width - Cm(0.4), Cm(0.65),
                nome, font_size=11, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Cm(0.2), top + Cm(0.65),
                width - Cm(0.4), Cm(0.4),
                url, font_size=8, bold=False,
                color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
    # Imagem ou área de texto
    if img_path:
        img_h = height - header_h - Cm(0.15)
        add_image_safe(slide, img_path, left, top + header_h, width, img_h, nome)
    else:
        body_top = top + header_h + Cm(0.2)
        body_h = height - header_h - Cm(0.3)
        body_box = slide.shapes.add_textbox(left + Cm(0.2), body_top,
                                            width - Cm(0.4), body_h)
        body_box.text_frame.word_wrap = True
        first = True
        for linha in descricao.split('\n'):
            if first:
                p = body_box.text_frame.paragraphs[0]
                first = False
            else:
                p = body_box.text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = "• " + linha
            run.font.size = Pt(10)
            run.font.name = FONTE_CORPO
            run.font.color.rgb = COR_CINZA_TEXTO


def new_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    setup_slide_background(slide, COR_BRANCO)
    return slide


def std_slide(prs, title, subtitle=None, subtitle_text=None):
    """Cria slide padrão com header, logo e rodapé."""
    slide = new_slide(prs)
    add_header_bar(slide, title, subtitle or subtitle_text)
    add_logo(slide)
    add_footer(slide)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# CRIAÇÃO DA APRESENTAÇÃO
# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

print("Criando 28 slides (v2)...")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════════
slide = new_slide(prs)
setup_slide_background(slide, COR_AZUL_ABAR)

# Faixa esquerda decorativa
add_rect(slide, Cm(0), Cm(0), Cm(0.6), SLIDE_H, COR_AZUL_CLARO)

# Painel branco central
panel_l = Cm(0.6)
panel_t = Cm(2.5)
panel_w = SLIDE_W - Cm(0.6) - Cm(3.5)
panel_h = SLIDE_H - Cm(4.5)
add_rect(slide, panel_l, panel_t, panel_w, panel_h, COR_BRANCO)

# Logo
if os.path.exists(IMG_LOGO):
    slide.shapes.add_picture(IMG_LOGO, Cm(1.2), Cm(0.4), Cm(6.5), Cm(2.0))

# Linha decorativa
add_rect(slide, panel_l, panel_t, panel_w, Cm(0.1), COR_AZUL_CLARO)

# Título principal
title_box = slide.shapes.add_textbox(Cm(1.2), Cm(3.2), Cm(25), Cm(4.0))
title_box.text_frame.word_wrap = True
p = title_box.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Bloco 6 — Auditoria de Dados,"
run.font.size = Pt(28)
run.font.bold = True
run.font.name = FONTE_TITULO
run.font.color.rgb = COR_AZUL_ABAR

add_paragraph(title_box.text_frame, "BI e Transparência Reguladora",
              font_size=28, bold=True, color=COR_AZUL_ABAR, align=PP_ALIGN.LEFT)

# Subtítulo
sub_box = slide.shapes.add_textbox(Cm(1.2), Cm(7.5), Cm(25), Cm(2.0))
sub_box.text_frame.word_wrap = True
p = sub_box.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Como a Inteligência Artificial transforma"
run.font.size = Pt(17)
run.font.name = FONTE_CORPO
run.font.color.rgb = COR_AZUL_CLARO
add_paragraph(sub_box.text_frame, "a ciência de dados para reguladores",
              font_size=17, color=COR_AZUL_CLARO, align=PP_ALIGN.LEFT)

# Linha divisória
add_rect(slide, Cm(1.2), Cm(9.8), Cm(22), Cm(0.06), COR_AZUL_CLARO)

# Professor
add_textbox(slide, Cm(1.2), Cm(10.2), Cm(22), Cm(0.9),
            "Prof. Raul Araújo — AGENERSA",
            font_size=15, bold=True, color=COR_AZUL_ABAR)

# Rodapé da capa
add_rect(slide, panel_l, SLIDE_H - Cm(1.5), panel_w, Cm(1.5), COR_AZUL_ABAR)
add_textbox(slide, Cm(1.2), SLIDE_H - Cm(1.4), Cm(22), Cm(1.2),
            "Curso ABAR — Medições Inteligentes e Gestão Integrada  |  Fevereiro 2026",
            font_size=11, color=COR_BRANCO, align=PP_ALIGN.LEFT)

# Imagem decorativa lateral direita
add_image_safe(slide, IMG_BG, SLIDE_W - Cm(8.5), Cm(0), Cm(8.5), SLIDE_H, "background")

print("  Slide 1 OK — Capa")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — A IA Evoluiu Mais em 2 Anos
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "A IA Evoluiu Mais em 2 Anos do Que em 20",
                  subtitle_text="De curiosidade acadêmica a ferramenta profissional indispensável")

bullets = [
    ("2022", "IA gerava texto confuso, errava contas simples — considerada brinquedo", COR_CINZA_TEXTO),
    ("2024", "GPT-4 aprovado no Bar Exam (top 10%) — nível de advogado", COR_CINZA_TEXTO),
    ("2025", "Medalha de Ouro na Olimpíada Internacional de Matemática (IMO) — inédito", COR_AZUL_ABAR),
    ("2025", "SWE-bench: resolve 76,8% de bugs reais em repositórios GitHub", COR_AZUL_ABAR),
    ("2025", "Humanity's Last Exam: 90% de acertos em questões de PhDs", COR_AZUL_ABAR),
    ("Reflexão", "\"Se uma IA ganha medalha de ouro em matemática... o que pode fazer com seus dados?\"", COR_VERDE_DESTAQUE),
]

y = CONTENT_TOP + Cm(0.3)
row_h = Cm(1.6)
for ano, texto, txt_color in bullets:
    badge_color = COR_AZUL_ABAR if ano not in ("Reflexão",) else COR_VERDE_DESTAQUE
    add_rect(slide, CONTENT_L, y, Cm(2.2), Cm(0.8), badge_color)
    add_textbox(slide, CONTENT_L, y, Cm(2.2), Cm(0.8),
                ano, font_size=9, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    txt = slide.shapes.add_textbox(Cm(3.0), y, Cm(16.5), Cm(0.85))
    txt.text_frame.word_wrap = True
    p = txt.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(11.5)
    run.font.name = FONTE_CORPO
    run.font.color.rgb = txt_color
    run.font.bold = (ano == "Reflexão")
    run.font.italic = (ano == "Reflexão")
    y += row_h

# Grid 2x2 de benchmarks
bm_l = Cm(20.5)
bm_t = CONTENT_TOP
bm_w = Cm(6.3)
bm_h = Cm(4.5)
gap  = Cm(0.3)
benchmarks = [IMG_IMO, IMG_SWEBENCH, IMG_EPOCH, IMG_HLE]
labels = ["IMO 2025", "SWE-bench", "Epoch AI", "Humanity's Last Exam"]
for i, (bm_img, lbl) in enumerate(zip(benchmarks, labels)):
    col = i % 2
    row = i // 2
    left = bm_l + col * (bm_w + gap)
    top  = bm_t + row * (bm_h + gap)
    add_image_safe(slide, bm_img, left, top, bm_w, bm_h, lbl)

print("  Slide 2 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Roteiro das 4 Horas
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "O Que Faremos Hoje — Roteiro das 4 Horas",
                  subtitle_text="Fio condutor: IBM Data Science Methodology aplicada com IA")

modulos = [
    ("Módulo 1", "30 min", "Ecossistema de ferramentas IA\n4 modelos + IDEs + CLIs + Agentes\nNíveis: do dia a dia ao projeto completo", COR_AZUL_ABAR),
    ("Módulo 2", "60 min", "Ciclo IBM Data Science com IA\nCada etapa acelerada por IA\nDados reais: gov.br e Kaggle", COR_AZUL_CLARO),
    ("Módulo 3", "50 min", "Demo ao vivo — do Excel ao relatório\n7 Notebooks + 28 chamadas Gemini\nRelatório 8,9 MB em 1 min 42 seg", COR_VERDE_DESTAQUE),
    ("Módulo 4", "15 min", "NotebookLM + Dashboards\nApostila técnica no chat\nObservatório de Gás e ANP", COR_ROXO),
]

block_w = CONTENT_W / 4 - Cm(0.4)
block_h = CONTENT_H - Cm(1.2)
y_start = CONTENT_TOP + Cm(0.3)

for i, (titulo, duracao, corpo, cor) in enumerate(modulos):
    x = CONTENT_L + i * (block_w + Cm(0.4))
    add_rect(slide, x, y_start, block_w, block_h, COR_CINZA_CLARO, cor, Pt(2))
    add_rect(slide, x, y_start, block_w, Cm(1.9), cor)
    add_textbox(slide, x + Cm(0.2), y_start + Cm(0.1),
                block_w - Cm(0.4), Cm(0.7),
                titulo, font_size=13, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Cm(0.2), y_start + Cm(0.85),
                block_w - Cm(0.4), Cm(0.9),
                duracao, font_size=17, bold=True, color=COR_AMARELO_DEMO, align=PP_ALIGN.CENTER)
    y_txt = y_start + Cm(2.1)
    for linha in corpo.split('\n'):
        tb = slide.shapes.add_textbox(x + Cm(0.3), y_txt, block_w - Cm(0.6), Cm(0.9))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "• " + linha
        run.font.size = Pt(10)
        run.font.name = FONTE_CORPO
        run.font.color.rgb = COR_CINZA_TEXTO
        y_txt += Cm(1.05)

# Barra de intervalo
bar_top = y_start + block_h + Cm(0.2)
add_rect(slide, CONTENT_L, bar_top, CONTENT_W, Cm(0.65), COR_CINZA_MEDIO)
add_textbox(slide, CONTENT_L, bar_top + Cm(0.05), CONTENT_W, Cm(0.55),
            "⏸  Intervalo de 10 minutos entre Módulos 2 e 3",
            font_size=10, color=COR_AZUL_ABAR, align=PP_ALIGN.CENTER)

print("  Slide 3 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Nível 1: Os 4 Grandes Modelos (NOVO)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Nível 1 — Os 4 Grandes Modelos de IA",
                  subtitle_text="Acesso pelo navegador, sem instalar nada — gratuitos para começar")

# Grid 2×2 com screenshots + nome do modelo
grid_imgs = [
    (IMG_CHATGPT, "ChatGPT", "OpenAI", "chatgpt.com", COR_VERDE_DESTAQUE),
    (IMG_GEMINI,  "Gemini",  "Google", "gemini.google.com", COR_AZUL_CLARO),
    (IMG_CLAUDE,  "Claude",  "Anthropic", "claude.ai", COR_LARANJA),
    (IMG_GROK,    "Grok",    "xAI (Elon Musk)", "x.com/i/grok", COR_CINZA_TEXTO),
]

cell_w = CONTENT_W / 2 - Cm(0.4)
cell_h = CONTENT_H / 2 - Cm(0.3)

for i, (img_path, nome, empresa, url, cor) in enumerate(grid_imgs):
    col = i % 2
    row = i // 2
    left = CONTENT_L + col * (cell_w + Cm(0.4))
    top  = CONTENT_TOP + Cm(0.2) + row * (cell_h + Cm(0.3))

    # Borda do card
    add_rect(slide, left, top, cell_w, cell_h, COR_CINZA_CLARO, cor, Pt(2))

    # Header do card
    header_h = Cm(0.9)
    add_rect(slide, left, top, cell_w, header_h, cor)
    add_textbox(slide, left + Cm(0.2), top + Cm(0.05),
                cell_w - Cm(0.4), Cm(0.5),
                f"{nome}  ·  {empresa}  ·  {url}",
                font_size=9, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

    # Screenshot
    img_area_h = cell_h - header_h - Cm(0.1)
    add_image_safe(slide, img_path, left, top + header_h, cell_w, img_area_h, nome)

print("  Slide 4 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Nível 1: Capacidades Avançadas (NOVO)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Nível 1 — Além do Chat: Capacidades Avançadas",
                  subtitle_text="Ferramentas que vão muito além de responder perguntas")

capacidades = [
    ("Deep Research", COR_VERDE_DESTAQUE,
     "Pesquisa autônoma na web\nAnálise de dezenas de fontes em minutos\nGera relatório estruturado com referências",
     IMG_FRONTIER),
    ("Canvas / Artifacts", COR_AZUL_ABAR,
     "Lousa colaborativa em tempo real\nEdição de documentos, código, apresentações\nClaude Artifacts: componentes interativos",
     IMG_CHATBOT),
    ("Code Interpreter / ADA", COR_LARANJA,
     "Python rodando dentro do chat\nGráficos, análise de planilhas, limpeza de dados\nCarregue seu Excel → gráficos em segundos",
     IMG_ARTIFICIAL),
]

col_w = CONTENT_W / 3 - Cm(0.3)
col_h = CONTENT_H - Cm(0.2)

for i, (nome, cor, descr, bm_img) in enumerate(capacidades):
    left = CONTENT_L + i * (col_w + Cm(0.3))
    top  = CONTENT_TOP + Cm(0.2)

    add_rect(slide, left, top, col_w, col_h, COR_CINZA_CLARO, cor, Pt(2))
    # Header
    add_rect(slide, left, top, col_w, Cm(0.9), cor)
    add_textbox(slide, left + Cm(0.2), top + Cm(0.1),
                col_w - Cm(0.4), Cm(0.7),
                nome, font_size=12, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

    # Screenshot de benchmark como ilustração
    img_h = Cm(6.5)
    add_image_safe(slide, bm_img, left, top + Cm(0.9), col_w, img_h, nome)

    # Bullets
    y_txt = top + Cm(0.9) + img_h + Cm(0.2)
    for linha in descr.split('\n'):
        tb = slide.shapes.add_textbox(left + Cm(0.2), y_txt, col_w - Cm(0.4), Cm(0.8))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "• " + linha
        run.font.size = Pt(10)
        run.font.name = FONTE_CORPO
        run.font.color.rgb = COR_CINZA_TEXTO
        y_txt += Cm(0.85)

print("  Slide 5 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Nível 2a: IDEs com IA
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Nível 2a — IDEs com IA Integrada",
                  subtitle_text="Para desenvolvedores e analistas — IA dentro do ambiente de código")

ides = [
    ("VS Code + GitHub Copilot", "code.visualstudio.com",
     "Editor open-source da Microsoft\nCopilot integrado: sugestões, chat, completar código\nGratuito para estudantes e open-source",
     COR_AZUL_ABAR, IMG_VSCODE),
    ("Cursor AI", "cursor.com",
     "Fork do VS Code otimizado para IA\nCompose: IA cria arquivos inteiros\nMulti-file edits, terminal IA integrado",
     COR_AZUL_CLARO, IMG_CURSOR),
    ("Windsurf (ex-Codeium)", "windsurf.com",
     "IDE criado pela equipe ex-Codeium\nFlows: agentes que editam múltiplos arquivos\nAcquirida pelo Google em 2025",
     COR_VERDE_DESTAQUE, IMG_WINDSURF),
]

col_w = CONTENT_W / 3 - Cm(0.3)
col_h = CONTENT_H - Cm(0.3)

for i, (nome, url, descr, cor, img_path) in enumerate(ides):
    left = CONTENT_L + i * (col_w + Cm(0.3))
    top  = CONTENT_TOP + Cm(0.15)
    add_card_tool(slide, left, top, col_w, col_h, nome, url, descr, cor, img_path)

print("  Slide 6 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Nível 2b: CLIs com IA (NOVO)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Nível 2b — CLIs com IA para Projetos Completos",
                  subtitle_text="Rodam no terminal — IA acessa arquivos, executa código e itera sozinha")

clis = [
    ("Claude Code", "claude.ai/code",
     "CLI da Anthropic para projetos inteiros\nLê, edita e executa arquivos localmente\nAgentes e subagentes em paralelo\nEsta apresentação foi feita com Claude Code",
     COR_LARANJA, None),
    ("OpenAI Codex CLI", "platform.openai.com",
     "Baseado em GPT-4.1 otimizado para código\nRoda no terminal, acessa o sistema de arquivos\nComandos Unix integrados\nPara automações de DevOps",
     COR_VERDE_DESTAQUE, IMG_CODEX),
    ("Gemini CLI", "ai.google.dev/gemini-api",
     "CLI oficial do Google Gemini\nContexto de 1 milhão de tokens\nIdeal para projetos com muitos arquivos\nIntegrado ao Google Cloud",
     COR_AZUL_CLARO, None),
]

col_w = CONTENT_W / 3 - Cm(0.3)
col_h = CONTENT_H - Cm(0.3)

for i, (nome, url, descr, cor, img_path) in enumerate(clis):
    left = CONTENT_L + i * (col_w + Cm(0.3))
    top  = CONTENT_TOP + Cm(0.15)
    add_card_tool(slide, left, top, col_w, col_h, nome, url, descr, cor, img_path)

# Destaque no card do Claude Code
add_rect(slide, CONTENT_L + Cm(0.1), CONTENT_TOP + Cm(0.15) + Cm(10.5),
         col_w - Cm(0.2), Cm(0.8), COR_AMARELO_DEMO)
add_textbox(slide, CONTENT_L + Cm(0.2), CONTENT_TOP + Cm(0.15) + Cm(10.5),
            col_w - Cm(0.4), Cm(0.8),
            "★ Ferramenta usada para criar esta apresentação",
            font_size=9, bold=True, color=RGBColor(0x80, 0x30, 0x00),
            align=PP_ALIGN.CENTER)

print("  Slide 7 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Nível 3: Agentes Autônomos
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Nível 3 — Agentes Autônomos",
                  subtitle_text="IA que executa tarefas complexas em múltiplos passos sem intervenção humana")

agentes = [
    ("Claude Cowork", "anthropic.com",
     "Agente colaborativo da Anthropic\nMultiplos sub-agentes trabalhando juntos\nResearch preview — futuro da IA no trabalho\nDivide tarefas, delega, consolida resultado",
     COR_LARANJA, IMG_COWORK),
    ("Manus", "manus.im",
     "Agente autônomo de tarefas gerais\nPlaneja, pesquisa, codifica e entrega\nUsa ferramentas: browser, código, arquivos\nResultado: produto final, não só resposta",
     COR_AZUL_ABAR, IMG_MANUS),
    ("Genspark", "genspark.ai",
     "Agente de pesquisa profunda\nGera Sparks: mini-apps interativos\nAgrega múltiplas fontes com citações\nIdeal para análise de mercado e regulação",
     COR_AZUL_CLARO, IMG_GENSPARK),
]

col_w = CONTENT_W / 3 - Cm(0.3)
col_h = CONTENT_H - Cm(0.3)

for i, (nome, url, descr, cor, img_path) in enumerate(agentes):
    left = CONTENT_L + i * (col_w + Cm(0.3))
    top  = CONTENT_TOP + Cm(0.15)
    add_card_tool(slide, left, top, col_w, col_h, nome, url, descr, cor, img_path)

# Mensagem de diferencial
diff_top = SLIDE_H - FOOTER_H - Cm(1.1)
add_rect(slide, CONTENT_L, diff_top, CONTENT_W, Cm(0.8), COR_AZUL_ABAR)
add_textbox(slide, CONTENT_L + Cm(0.3), diff_top + Cm(0.1), CONTENT_W - Cm(0.6), Cm(0.65),
            "Diferença fundamental: Você dá a META — o agente define o PLANO e EXECUTA",
            font_size=10, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

print("  Slide 8 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Escolhendo por Tipo de Tarefa
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Escolhendo a Ferramenta por Tipo de Tarefa",
                  subtitle_text="A ferramenta certa depende do que você quer fazer, não de quão avançado você é")

tarefas = [
    ("Redação e Relatórios",    "Claude 4.6 / ChatGPT 5.2",  "Relatórios, contratos, análises, atas",      COR_AZUL_ABAR),
    ("Código e Automação",      "Claude Code / Codex CLI",   "Scripts Python, pipelines, notebooks",       COR_AZUL_CLARO),
    ("Análise de Planilhas",    "ChatGPT ADA / Claude Web",  "Excel, CSV, EDA, limpeza de dados",          COR_VERDE_DESTAQUE),
    ("Pesquisa Aprofundada",    "Gemini Deep Research / Genspark", "Artigos, relatórios setoriais, normas", COR_LARANJA),
    ("Imagens e Diagramas",     "DALL-E / Midjourney",       "Infográficos, slides, logos, diagramas",     COR_ROXO),
    ("Projetos Completos (end-to-end)", "Agentes: Manus / Cowork", "Da pergunta ao produto final",        COR_VERMELHO),
]

# Tabela visual com 2 colunas
half = len(tarefas) // 2
row_h = Cm(2.0)
col_w = CONTENT_W / 2 - Cm(0.3)
y0 = CONTENT_TOP + Cm(0.2)

for i, (tarefa, ferramenta, exemplos, cor) in enumerate(tarefas):
    col = i // half
    row = i % half
    left = CONTENT_L + col * (col_w + Cm(0.3))
    top  = y0 + row * (row_h + Cm(0.2))

    # Fundo da linha
    add_rect(slide, left, top, col_w, row_h, COR_CINZA_CLARO, cor, Pt(1))
    # Badge cor lateral
    add_rect(slide, left, top, Cm(0.4), row_h, cor)

    # Texto tarefa
    add_textbox(slide, left + Cm(0.6), top + Cm(0.05),
                col_w - Cm(0.7), Cm(0.65),
                tarefa, font_size=11, bold=True, color=COR_AZUL_ABAR)
    # Ferramenta
    add_textbox(slide, left + Cm(0.6), top + Cm(0.65),
                col_w - Cm(0.7), Cm(0.55),
                "→  " + ferramenta, font_size=10, bold=False, color=cor)
    # Exemplos
    add_textbox(slide, left + Cm(0.6), top + Cm(1.2),
                col_w - Cm(0.7), Cm(0.65),
                exemplos, font_size=9, color=COR_CINZA_TEXTO)

print("  Slide 9 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Demo 1a: ChatGPT + Excel
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Demo 1a — ChatGPT: Análise de Planilha Excel",
                  subtitle_text="Sem instalar nada, sem escrever código — só copiar e colar")

add_demo_banner(slide, CONTENT_TOP)

y_after_banner = CONTENT_TOP + Cm(1.1)

# Coluna esquerda: passos
passos = [
    ("1", "Abrir chatgpt.com no navegador"),
    ("2", "Fazer upload da planilha de dados do distrito (xlsx)"),
    ("3", "Prompt: 'Faça uma análise exploratória dos dados. Mostre: número de linhas, colunas, valores nulos, estatísticas descritivas e gere um gráfico de boxplot das pressões por medidor'"),
    ("4", "Observar: ChatGPT executa Python internamente e retorna gráfico + análise"),
    ("5", "Pedir variações: 'Agora compare os quartis de vazão entre meses'"),
]
y = y_after_banner + Cm(0.3)
for num, texto in passos:
    add_rect(slide, CONTENT_L, y, Cm(0.7), Cm(0.65), COR_AZUL_ABAR)
    add_textbox(slide, CONTENT_L, y, Cm(0.7), Cm(0.65),
                num, font_size=11, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Cm(1.4), y, Cm(16.5), Cm(0.75))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(11)
    run.font.name = FONTE_CORPO
    run.font.color.rgb = COR_CINZA_TEXTO
    y += Cm(1.1)

# Imagem lateral
add_image_safe(slide, IMG_CHATGPT,
               Cm(20.0), y_after_banner,
               SLIDE_W - Cm(20.5), CONTENT_H - Cm(1.2), "ChatGPT Interface")

print("  Slide 10 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Demo 1b: Claude + Excel (NOVO)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Demo 1b — Claude: Relatório Técnico a partir do Excel",
                  subtitle_text="Claude vai além do gráfico: entende contexto técnico e escreve análises completas")

add_demo_banner(slide, CONTENT_TOP)

y_after_banner = CONTENT_TOP + Cm(1.1)

# Comparação ChatGPT vs Claude
comp_col_w = CONTENT_W / 2 - Cm(0.3)

# ChatGPT lado
add_rect(slide, CONTENT_L, y_after_banner, comp_col_w, Cm(1.2), COR_VERDE_DESTAQUE)
add_textbox(slide, CONTENT_L + Cm(0.2), y_after_banner + Cm(0.2),
            comp_col_w - Cm(0.4), Cm(0.8),
            "ChatGPT ADA — Ideal para exploração", font_size=11, bold=True, color=COR_BRANCO)

pontos_gpt = ["Gera gráficos interativos", "Análise rápida e visual", "Excelente para EDA inicial",
              "Executa Python dentro do chat"]
y = y_after_banner + Cm(1.35)
for ponto in pontos_gpt:
    add_textbox(slide, CONTENT_L + Cm(0.3), y, comp_col_w - Cm(0.5), Cm(0.65),
                "✓  " + ponto, font_size=10, color=COR_CINZA_TEXTO)
    y += Cm(0.75)

# Claude lado
claude_l = CONTENT_L + comp_col_w + Cm(0.3)
add_rect(slide, claude_l, y_after_banner, comp_col_w, Cm(1.2), COR_LARANJA)
add_textbox(slide, claude_l + Cm(0.2), y_after_banner + Cm(0.2),
            comp_col_w - Cm(0.4), Cm(0.8),
            "Claude — Ideal para relatórios técnicos", font_size=11, bold=True, color=COR_BRANCO)

pontos_claude = ["Interpreta contexto técnico profundo", "Escreve análises com metodologia",
                 "Referencia normas (GUM, ABNT, ANP)", "Ideal para relatórios regulatórios"]
y = y_after_banner + Cm(1.35)
for ponto in pontos_claude:
    add_textbox(slide, claude_l + Cm(0.3), y, comp_col_w - Cm(0.5), Cm(0.65),
                "✓  " + ponto, font_size=10, color=COR_CINZA_TEXTO)
    y += Cm(0.75)

# Imagem do Claude Interface
img_top = y_after_banner + Cm(5.5)
img_h = CONTENT_H - Cm(1.2) - Cm(5.5)
add_image_safe(slide, IMG_CLAUDE, CONTENT_L, img_top, CONTENT_W, img_h, "Claude Interface")

print("  Slide 11 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — IBM Cycle + IA (numerado e expandido)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "O Ciclo IBM Data Science — Potencializado por IA",
                  subtitle_text="Cada etapa da metodologia clássica tem uma ferramenta IA ideal")

# Diagrama IBM à esquerda
img_w = Cm(14.5)
img_h = CONTENT_H - Cm(0.2)
add_image_safe(slide, IMG_IBM, CONTENT_L, CONTENT_TOP + Cm(0.1), img_w, img_h, "Ciclo IBM")

# Lista de etapas à direita
etapas = [
    ("1", "Entend. Negócio",   "NotebookLM + Claude multimodal"),
    ("2", "Abordagem Analítica","Claude (raciocínio), o1 (matemática)"),
    ("3", "Requisitos de Dados","ChatGPT, Gemini Deep Research"),
    ("4", "Coleta dos Dados",  "Firecrawl, Claude Code, gov.br"),
    ("5", "Entend. dos Dados", "ChatGPT ADA, Code Interpreter"),
    ("6", "Preparação",        "Claude Code, Copilot, pandas+IA"),
    ("7", "Criação do Modelo", "AutoML+GPT, Data Interpreter"),
    ("8", "Avaliação",         "ChatGPT, Claude (interpretação)"),
    ("9", "Deployment",        "FastAPI + Gemini (projeto demo)"),
]

right_l = CONTENT_L + img_w + Cm(0.5)
right_w = CONTENT_W - img_w - Cm(0.5)
row_h = CONTENT_H / len(etapas)
y = CONTENT_TOP + Cm(0.1)

for num, etapa, ia in etapas:
    add_rect(slide, right_l, y, Cm(0.6), row_h - Cm(0.08), COR_AZUL_ABAR)
    add_textbox(slide, right_l, y, Cm(0.6), row_h - Cm(0.08),
                num, font_size=8, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, right_l + Cm(0.7), y, right_w - Cm(0.75), Cm(0.65),
                etapa, font_size=9, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, right_l + Cm(0.7), y + Cm(0.6), right_w - Cm(0.75), Cm(0.55),
                "↳ " + ia, font_size=8.5, color=COR_CINZA_TEXTO)
    y += row_h

print("  Slide 12 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Etapa 1: Multimodalidade (Entendimento do Negócio)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Etapa 1 — Entendimento do Negócio com IA Multimodal",
                  subtitle_text="IA que enxerga documentos, imagens, planilhas e vídeos — não só texto")

# Três cards de capacidades multimodais
caps = [
    ("📄  Documentos PDF", COR_AZUL_ABAR,
     "Carregue a apostila do Prof. Barateiro no NotebookLM\n"
     "Pergunte: 'Qual a equação de incerteza Tipo A?'\n"
     "IA responde com referência à página exata\n"
     "Ideal para normas, manuais, regulamentos"),
    ("📊  Planilhas e Tabelas", COR_VERDE_DESTAQUE,
     "Envie o Excel de dados do distrito\n"
     "Prompt: 'Identifique anomalias nesta tabela'\n"
     "IA analisa automaticamente sem código\n"
     "Funciona com ChatGPT, Claude e Gemini"),
    ("🖼️  Imagens e Fotos", COR_LARANJA,
     "Foto de medidor, documento físico, diagrama\n"
     "Prompt: 'Leia o valor neste medidor de gás'\n"
     "Ou: 'Identifique problemas neste diagrama'\n"
     "GPT-4o e Claude 4.6 são excelentes nisso"),
]

col_w = CONTENT_W / 3 - Cm(0.3)
col_h = CONTENT_H - Cm(0.2)

for i, (titulo, cor, descr) in enumerate(caps):
    left = CONTENT_L + i * (col_w + Cm(0.3))
    top  = CONTENT_TOP + Cm(0.1)
    add_rect(slide, left, top, col_w, col_h, COR_CINZA_CLARO, cor, Pt(2))
    add_rect(slide, left, top, col_w, Cm(1.0), cor)
    add_textbox(slide, left + Cm(0.2), top + Cm(0.15),
                col_w - Cm(0.4), Cm(0.7),
                titulo, font_size=11, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    y_txt = top + Cm(1.2)
    for linha in descr.split('\n'):
        tb = slide.shapes.add_textbox(left + Cm(0.3), y_txt, col_w - Cm(0.5), Cm(0.9))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "• " + linha
        run.font.size = Pt(10.5)
        run.font.name = FONTE_CORPO
        run.font.color.rgb = COR_CINZA_TEXTO
        y_txt += Cm(0.95)

print("  Slide 13 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Etapa 2: Dados Abertos (gov.br / Kaggle)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Etapa 2 — Fontes de Dados Abertos para Reguladores",
                  subtitle_text="Milhares de datasets prontos para análise — gratuitos e públicos")

# Lado esquerdo: Kaggle
half_w = CONTENT_W / 2 - Cm(0.3)
top_img = CONTENT_TOP + Cm(0.2)
img_h = Cm(8.0)
txt_top = top_img + img_h + Cm(0.2)

add_image_safe(slide, IMG_KAGGLE, CONTENT_L, top_img, half_w, img_h, "Kaggle Housing Prices")
add_rect(slide, CONTENT_L, txt_top, half_w, Cm(1.1), COR_AZUL_ABAR)
add_textbox(slide, CONTENT_L + Cm(0.2), txt_top + Cm(0.1),
            half_w - Cm(0.4), Cm(0.9),
            "kaggle.com — California Housing Prices", font_size=10, bold=True, color=COR_BRANCO)
for i, linha in enumerate([
    "20.640 propriedades, 9 variáveis",
    "Benchmark clássico de análise de dados",
    "Usaremos para demo de limpeza e EDA",
]):
    add_textbox(slide, CONTENT_L + Cm(0.3), txt_top + Cm(1.2) + i * Cm(0.75),
                half_w - Cm(0.5), Cm(0.65),
                "• " + linha, font_size=10, color=COR_CINZA_TEXTO)

# Lado direito: dados.gov.br
right_l = CONTENT_L + half_w + Cm(0.3)
add_image_safe(slide, IMG_GOVBR, right_l, top_img, half_w, img_h, "dados.gov.br")
add_rect(slide, right_l, txt_top, half_w, Cm(1.1), COR_VERDE_DESTAQUE)
add_textbox(slide, right_l + Cm(0.2), txt_top + Cm(0.1),
            half_w - Cm(0.4), Cm(0.9),
            "dados.gov.br — Portal de Dados Abertos", font_size=10, bold=True, color=COR_BRANCO)
for i, linha in enumerate([
    "Dados de concessionárias, ANEEL, ANP, IBGE",
    "Formato aberto: CSV, JSON, API",
    "Para reguladores: dados já estão aqui",
]):
    add_textbox(slide, right_l + Cm(0.3), txt_top + Cm(1.2) + i * Cm(0.75),
                half_w - Cm(0.5), Cm(0.65),
                "• " + linha, font_size=10, color=COR_CINZA_TEXTO)

print("  Slide 14 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Etapa 2: Web Scraping com IA
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Etapa 2 — Coleta: Web Scraping com IA",
                  subtitle_text="Dados que não estão em CSV — IA raspa, estrutura e entrega")

# Coluna esquerda: como funciona
bullets = [
    ("O Problema", "Dados da ANP, ANEEL e AGENERSA publicados como tabelas em páginas HTML", COR_AZUL_ABAR),
    ("A Solução", "Pedir ao Claude Code: 'Raspe os dados de produção de gás do site da ANP dos últimos 12 meses e salve em CSV'", COR_AZUL_CLARO),
    ("O Resultado", "Claude escreve e executa o script Python com BeautifulSoup ou Playwright, entrega o CSV", COR_VERDE_DESTAQUE),
    ("Ferramentas", "Claude Code + Firecrawl (scraping avançado) + requests + pandas", COR_LARANJA),
]

y = CONTENT_TOP + Cm(0.3)
left_w = Cm(17.0)
for titulo, texto, cor in bullets:
    add_rect(slide, CONTENT_L, y, Cm(1.2), Cm(1.8), cor)
    add_textbox(slide, CONTENT_L, y + Cm(0.5), Cm(1.2), Cm(0.9),
                titulo[:5], font_size=8, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Cm(1.9), y + Cm(0.2), left_w - Cm(1.9) - Cm(0.3), Cm(1.5))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"{titulo}: {texto}"
    run.font.size = Pt(11)
    run.font.name = FONTE_CORPO
    run.font.color.rgb = COR_CINZA_TEXTO
    y += Cm(2.1)

# Imagem direita
add_image_safe(slide, IMG_GOVBR, Cm(18.0), CONTENT_TOP,
               SLIDE_W - Cm(18.5), CONTENT_H, "Dados.gov.br")

print("  Slide 15 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Etapas 3-4: Limpeza de Dados com Kaggle Housing
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Etapas 3-4 — Entendimento e Limpeza de Dados",
                  subtitle_text="Exemplo prático: California Housing Prices do Kaggle")

# Preview dos dados (texto simulando output do pandas)
code_top = CONTENT_TOP + Cm(0.2)
code_h = Cm(5.5)
add_rect(slide, CONTENT_L, code_top, Cm(19.0), code_h,
         RGBColor(0x1E, 0x1E, 0x2E))  # fundo escuro tipo terminal

code_lines = [
    ">>> df.shape",
    "(20640, 9)",
    ">>> df.isnull().sum()",
    "total_bedrooms    207   # ← valores nulos!",
    ">>> df.describe()",
    "       longitude   latitude  housing_median_age  ...",
    "mean    -119.57      35.63        28.6           ...",
    "std        2.00       2.14        12.6           ...",
]
y_code = code_top + Cm(0.25)
for line in code_lines:
    add_textbox(slide, CONTENT_L + Cm(0.3), y_code, Cm(18.4), Cm(0.55),
                line, font_size=8.5, color=RGBColor(0xA8, 0xFF, 0xA0),
                font_name="Courier New")
    y_code += Cm(0.6)

# Prompts usados
prompts_top = code_top + code_h + Cm(0.3)
prompts = [
    "Prompt 1: 'Analise este CSV. Quantos valores nulos existem? Onde?'",
    "Prompt 2: 'Preencha os valores nulos de total_bedrooms com a mediana por housing_median_age'",
    "Prompt 3: 'Crie um mapa de calor com a correlação entre as variáveis numéricas'",
]
for i, prompt in enumerate(prompts):
    add_rect(slide, CONTENT_L, prompts_top + i * Cm(1.1), Cm(19.0), Cm(0.9),
             COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(1))
    add_textbox(slide, CONTENT_L + Cm(0.3), prompts_top + i * Cm(1.1) + Cm(0.15),
                Cm(18.4), Cm(0.65), prompt, font_size=10, color=COR_AZUL_ABAR)

# Imagem do Kaggle
add_image_safe(slide, IMG_KAGGLE,
               Cm(20.0), CONTENT_TOP,
               SLIDE_W - Cm(20.5), CONTENT_H, "Kaggle Housing")

print("  Slide 16 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Etapa 5: Criação do Modelo — Notebooks do Projeto
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Etapa 5 — Criação do Modelo: Os 7 Notebooks do Projeto",
                  subtitle_text="Cada notebook gerado com IA — 'Não escrevi uma linha de Python'")

notebooks = [
    ("01", "Leitura e Exploração",  "183 dias de dados, 14 abas Excel, estatísticas iniciais"),
    ("02", "Análise de Pressão",    "Boxplots, outliers, histogramas por ponto de medição"),
    ("03", "Análise de Vazão",      "Séries temporais, sazonalidade, volumes diários"),
    ("04", "Qualidade do Gás",      "Parâmetros físico-químicos: densidade, poder calorífico"),
    ("05", "Eficiência Comercial",  "Perdas aparentes, balanço comercial por período"),
    ("06", "Incertezas (GUM)",      "Aplicação da metodologia GUM, propagação de incertezas"),
    ("07", "Balanço de Massa",      "Equilíbrio entrada-saída, análise de discrepâncias"),
]

row_h = CONTENT_H / len(notebooks) - Cm(0.1)
y = CONTENT_TOP + Cm(0.15)

for num, titulo, descricao in notebooks:
    add_rect(slide, CONTENT_L, y, Cm(1.5), row_h - Cm(0.05), COR_AZUL_ABAR)
    add_textbox(slide, CONTENT_L, y + (row_h - Cm(0.8)) / 2, Cm(1.5), Cm(0.8),
                num, font_size=14, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, Cm(1.6), y, CONTENT_W - Cm(1.6), row_h - Cm(0.05),
             COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(1))
    add_textbox(slide, Cm(1.9), y + Cm(0.1),
                Cm(11), row_h - Cm(0.2),
                titulo, font_size=11, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, Cm(13.0), y + Cm(0.1),
                CONTENT_W - Cm(13.5), row_h - Cm(0.2),
                descricao, font_size=10, color=COR_CINZA_TEXTO)
    y += row_h + Cm(0.1)

print("  Slide 17 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Demo 2: Os 7 Notebooks ao vivo
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Demo 2 — Os 7 Notebooks: Do Zero ao Relatório",
                  subtitle_text="Mostrar o projeto completo — como Claude Code gerou cada notebook")

add_demo_banner(slide, CONTENT_TOP)

y = CONTENT_TOP + Cm(1.1)

# Destaque central
add_rect(slide, CONTENT_L, y, CONTENT_W, Cm(2.8), COR_CINZA_CLARO, COR_AZUL_ABAR, Pt(2))
add_textbox(slide, CONTENT_L + Cm(0.5), y + Cm(0.3), CONTENT_W - Cm(1.0), Cm(2.2),
            "\"Abri um terminal, disse ao Claude Code: 'Tenho uma planilha com dados de um distrito de gás natural com 14 abas e 183 dias. "
            "Crie 7 notebooks Jupyter com análise completa, um para cada tema: pressão, vazão, qualidade do gás, eficiência comercial, "
            "incertezas GUM e balanço de massa.' Dois dias depois, o sistema estava pronto.\"",
            font_size=11, bold=False, color=COR_AZUL_ABAR)

y += Cm(3.1)

passos_demo = [
    ("Passo 1 (5 min)",  "Abrir 01_leitura_e_exploracao.ipynb — mostrar dados carregados e primeiros gráficos"),
    ("Passo 2 (5 min)",  "Rolar pelos demais notebooks — mostrar gráficos de pressão, vazão e qualidade"),
    ("Passo 3 (5 min)",  "Mostrar como o código foi gerado por IA (comentários, funções, docstrings)"),
]
for titulo, descr in passos_demo:
    add_rect(slide, CONTENT_L, y, Cm(3.5), Cm(1.1), COR_AZUL_ABAR)
    add_textbox(slide, CONTENT_L + Cm(0.1), y + Cm(0.2), Cm(3.3), Cm(0.7),
                titulo, font_size=9, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(Cm(4.2), y + Cm(0.15), CONTENT_W - Cm(4.3), Cm(0.85))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = descr
    run.font.size = Pt(11)
    run.font.name = FONTE_CORPO
    run.font.color.rgb = COR_CINZA_TEXTO
    y += Cm(1.4)

print("  Slide 18 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Arquitetura do Sistema
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Arquitetura do Sistema de Auditoria Automatizada",
                  subtitle_text="Excel → API Gemini (28 chamadas paralelas) → Relatório Word em 1 min 42 seg")

# Diagrama de fluxo textual à esquerda
etapas_arq = [
    ("📥 Entrada",        "Planilha Excel (2.1 MB)\n14 abas • 183 dias de dados"),
    ("⚙️  Processamento",  "7 Notebooks Python\nAnálise completa → 23 gráficos"),
    ("🤖 Gemini API",     "28 chamadas simultâneas\nGera texto para cada seção"),
    ("📄 Saída",           "Relatório Word (8.9 MB)\n100+ páginas em 1 min 42 seg"),
]

left_w = Cm(15.0)
box_h = CONTENT_H / len(etapas_arq) - Cm(0.4)
y = CONTENT_TOP + Cm(0.2)

for i, (icone_titulo, descr) in enumerate(etapas_arq):
    cor = [COR_AZUL_ABAR, COR_AZUL_CLARO, COR_VERDE_DESTAQUE, COR_LARANJA][i]
    add_rect(slide, CONTENT_L, y, left_w, box_h, COR_CINZA_CLARO, cor, Pt(2))
    add_rect(slide, CONTENT_L, y, left_w, Cm(1.0), cor)
    add_textbox(slide, CONTENT_L + Cm(0.3), y + Cm(0.15),
                left_w - Cm(0.6), Cm(0.75),
                icone_titulo, font_size=12, bold=True, color=COR_BRANCO)
    for j, linha in enumerate(descr.split('\n')):
        add_textbox(slide, CONTENT_L + Cm(0.4), y + Cm(1.2) + j * Cm(0.7),
                    left_w - Cm(0.7), Cm(0.65),
                    linha, font_size=10.5, color=COR_CINZA_TEXTO)
    if i < len(etapas_arq) - 1:
        arrow_y = y + box_h + Cm(0.1)
        add_textbox(slide, CONTENT_L + left_w / 2 - Cm(1), arrow_y - Cm(0.3),
                    Cm(2), Cm(0.5), "▼", font_size=14, bold=True,
                    color=cor, align=PP_ALIGN.CENTER)
    y += box_h + Cm(0.4)

# Diagrama à direita
add_image_safe(slide, IMG_FLUXO, Cm(16.5), CONTENT_TOP,
               SLIDE_W - Cm(17.0), CONTENT_H, "Fluxo de Auditoria")

print("  Slide 19 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Demo 3: Interface Web (grid 2×3 corrigido)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Demo 3 — Interface Web: Do Upload ao Download",
                  subtitle_text="FastAPI + Server-Sent Events — progresso em tempo real no navegador")

add_demo_banner(slide, CONTENT_TOP)

# Grid 2×3 de screenshots — CORRIGIDO com layout calculado
grid_imgs = [
    (IMG_WEB1, "1. Configuração"),
    (IMG_WEB2, "2. Pipeline em Progresso"),
    (IMG_WEB3, "3. Galeria de Gráficos"),
    (IMG_WEB4, "4. Diagramas IA"),
    (IMG_WEB5, "5. Textos Gerados"),
    (IMG_WEB6, "6. Downloads"),
]

n_cols = 3
n_rows = 2
banner_h = Cm(1.0)
avail_top = CONTENT_TOP + banner_h + Cm(0.2)
avail_h   = SLIDE_H - avail_top - FOOTER_H - Cm(0.2)
avail_w   = CONTENT_W
gap_x     = Cm(0.35)
gap_y     = Cm(0.3)
cell_w    = (avail_w - gap_x * (n_cols - 1)) / n_cols
cell_h    = (avail_h - gap_y * (n_rows - 1)) / n_rows
label_h   = Cm(0.5)
img_h     = cell_h - label_h

for i, (img_path, label) in enumerate(grid_imgs):
    col = i % n_cols
    row = i // n_cols
    left = CONTENT_L + col * (cell_w + gap_x)
    top  = avail_top + row * (cell_h + gap_y)
    # Label acima da imagem
    add_rect(slide, left, top, cell_w, label_h, COR_AZUL_ABAR)
    add_textbox(slide, left + Cm(0.1), top + Cm(0.05),
                cell_w - Cm(0.2), label_h - Cm(0.05),
                label, font_size=8, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    # Imagem
    add_image_safe(slide, img_path, left, top + label_h, cell_w, img_h, label)

print("  Slide 20 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — O Relatório Final
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "O Relatório Final — 100+ Páginas Geradas em < 2 Minutos",
                  subtitle_text="8,9 MB de relatório técnico: capítulos, equações, gráficos, apêndices")

# Stats em destaque
stats = [
    ("8,9 MB",    "Tamanho do arquivo Word"),
    ("100+",      "Páginas de relatório"),
    ("1:42 min",  "Tempo de geração"),
    ("28",        "Seções geradas em paralelo"),
    ("23",        "Gráficos inseridos"),
    ("7",         "Notebooks de análise"),
]

stat_w = CONTENT_W / 6 - Cm(0.2)
stat_h = Cm(2.8)
x = CONTENT_L
y = CONTENT_TOP + Cm(0.3)

for valor, descr in stats:
    add_rect(slide, x, y, stat_w, stat_h, COR_AZUL_ABAR)
    add_textbox(slide, x + Cm(0.1), y + Cm(0.2),
                stat_w - Cm(0.2), Cm(1.3),
                valor, font_size=20, bold=True, color=COR_AMARELO_DEMO,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Cm(0.1), y + Cm(1.5),
                stat_w - Cm(0.2), Cm(1.1),
                descr, font_size=9, color=COR_BRANCO,
                align=PP_ALIGN.CENTER)
    x += stat_w + Cm(0.2)

# Estrutura do relatório
struct_top = CONTENT_TOP + stat_h + Cm(0.7)
estrutura = [
    "Sumário Executivo",
    "1. Metodologia (GUM / ABNT NBR ISO 5725)",
    "2. Análise de Pressão — 8 medidores, 183 dias",
    "3. Análise de Vazão — volumes diários e mensais",
    "4. Qualidade do Gás — parâmetros físico-químicos",
    "5. Eficiência Comercial — perdas aparentes",
    "6. Incertezas de Medição — Tipo A e Tipo B",
    "7. Balanço de Massa — equilíbrio do distrito",
    "8. Conclusões e Recomendações",
    "Apêndice A — Código Python completo",
    "Apêndice B — Gráficos detalhados (23 figuras)",
]

left_col_w = Cm(17.0)
y = struct_top
for item in estrutura:
    add_textbox(slide, CONTENT_L + Cm(0.3), y, left_col_w - Cm(0.5), Cm(0.75),
                "▸ " + item, font_size=10, color=COR_CINZA_TEXTO)
    y += Cm(0.8)

add_image_safe(slide, IMG_PROCESSO, Cm(18.0), struct_top,
               SLIDE_W - Cm(18.5), CONTENT_H - stat_h - Cm(0.8), "Processo de Análise")

print("  Slide 21 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Comparativo Manual vs IA
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Comparativo — Análise Manual vs. IA Assistida",
                  subtitle_text="Mesmo dataset, mesma metodologia, resultados alinhados — tempo radicalmente diferente")

# Tabela comparativa
colunas = ["", "Análise Tradicional", "Com IA Assistida"]
linhas = [
    ("Tempo total",      "2–4 semanas",        "< 2 minutos"),
    ("Linguagem",        "Python/R necessário", "Português natural"),
    ("Código escrito",   "500–1000 linhas",     "Zero (gerado por IA)"),
    ("Gráficos",         "Manual, um a um",     "23 gráficos automáticos"),
    ("Relatório",        "Word manual",         "Word gerado, 8,9 MB"),
    ("Repetibilidade",   "Alta (se documentado)","Alta (código versionado)"),
    ("Qualidade técnica","Depende do analista",  "Metodologia GUM aplicada"),
]

col_w_lbl = Cm(8.0)
col_w_val = (CONTENT_W - col_w_lbl) / 2
row_h = (CONTENT_H - Cm(1.2)) / (len(linhas) + 1)
y = CONTENT_TOP + Cm(0.2)

# Header
for j, col_title in enumerate(colunas):
    col_l = CONTENT_L + (j * (col_w_val if j > 0 else col_w_lbl) +
                         (col_w_lbl if j > 1 else 0))
    if j == 0:
        col_l = CONTENT_L
        w = col_w_lbl
    elif j == 1:
        col_l = CONTENT_L + col_w_lbl
        w = col_w_val
    else:
        col_l = CONTENT_L + col_w_lbl + col_w_val
        w = col_w_val
    add_rect(slide, col_l, y, w, row_h - Cm(0.05), COR_AZUL_ABAR)
    add_textbox(slide, col_l + Cm(0.2), y + Cm(0.1), w - Cm(0.3), row_h - Cm(0.15),
                col_title, font_size=11, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
y += row_h

for k, (aspecto, manual, ia) in enumerate(linhas):
    bg = COR_CINZA_CLARO if k % 2 == 0 else COR_BRANCO
    # Aspecto
    add_rect(slide, CONTENT_L, y, col_w_lbl, row_h - Cm(0.05), bg)
    add_textbox(slide, CONTENT_L + Cm(0.2), y + Cm(0.05),
                col_w_lbl - Cm(0.3), row_h - Cm(0.15),
                aspecto, font_size=10, bold=True, color=COR_AZUL_ABAR)
    # Manual
    add_rect(slide, CONTENT_L + col_w_lbl, y, col_w_val, row_h - Cm(0.05), bg)
    add_textbox(slide, CONTENT_L + col_w_lbl + Cm(0.2), y + Cm(0.05),
                col_w_val - Cm(0.3), row_h - Cm(0.15),
                manual, font_size=10, color=COR_CINZA_TEXTO)
    # IA (verde claro de fundo)
    add_rect(slide, CONTENT_L + col_w_lbl + col_w_val, y, col_w_val, row_h - Cm(0.05),
             RGBColor(0xE8, 0xF8, 0xEE))
    add_textbox(slide, CONTENT_L + col_w_lbl + col_w_val + Cm(0.2), y + Cm(0.05),
                col_w_val - Cm(0.3), row_h - Cm(0.15),
                ia, font_size=10, bold=True, color=COR_VERDE_DESTAQUE)
    y += row_h

print("  Slide 22 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Dashboard: Observatório de Gás AGENERSA (NOVO)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Dashboard Real — Observatório de Gás Natural (AGENERSA)",
                  subtitle_text="Da análise Python ao dashboard público — dbagenersa.vercel.app")

# Screenshot grande do dashboard
img_w = Cm(22.0)
img_h = CONTENT_H - Cm(0.2)
add_image_safe(slide, IMG_OBS_GAS, CONTENT_L, CONTENT_TOP + Cm(0.1),
               img_w, img_h, "Observatório de Gás AGENERSA")

# Painel lateral com informações
panel_l = CONTENT_L + img_w + Cm(0.4)
panel_w = SLIDE_W - panel_l - Cm(0.3)
panel_h = CONTENT_H - Cm(0.2)
top = CONTENT_TOP + Cm(0.1)

add_rect(slide, panel_l, top, panel_w, panel_h, COR_CINZA_CLARO, COR_AZUL_ABAR, Pt(2))
add_rect(slide, panel_l, top, panel_w, Cm(1.0), COR_AZUL_ABAR)
add_textbox(slide, panel_l + Cm(0.2), top + Cm(0.15),
            panel_w - Cm(0.4), Cm(0.7),
            "Sobre o Projeto", font_size=11, bold=True,
            color=COR_BRANCO, align=PP_ALIGN.CENTER)

infos = [
    ("Link", "dbagenersa.vercel.app"),
    ("Stack", "Python + Dash + Vercel"),
    ("Dados", "AGENERSA (públicos)"),
    ("Parceria", "Prof. Beraldi (UFF)"),
    ("Deploy", "Gratuito — 0 custo"),
]
y_info = top + Cm(1.2)
for chave, valor in infos:
    add_textbox(slide, panel_l + Cm(0.2), y_info, panel_w - Cm(0.4), Cm(0.55),
                chave + ":", font_size=9, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, panel_l + Cm(0.2), y_info + Cm(0.5), panel_w - Cm(0.4), Cm(0.55),
                valor, font_size=9, color=COR_CINZA_TEXTO)
    y_info += Cm(1.2)

add_rect(slide, panel_l + Cm(0.2), y_info, panel_w - Cm(0.4), Cm(1.8),
         COR_AZUL_ABAR)
add_textbox(slide, panel_l + Cm(0.3), y_info + Cm(0.15), panel_w - Cm(0.6), Cm(1.5),
            "\"Python + Vercel = Dashboard público em horas, sem servidor, sem custo\"",
            font_size=9, bold=True, color=COR_BRANCO)

print("  Slide 23 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Dashboard: ANP Painéis Dinâmicos (NOVO)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Dashboard Regulatório — Painéis Dinâmicos da ANP",
                  subtitle_text="Agência Nacional do Petróleo, Gás Natural e Biocombustíveis")

# Screenshot grande
img_w = Cm(22.0)
img_h = CONTENT_H - Cm(0.2)
add_image_safe(slide, IMG_ANP, CONTENT_L, CONTENT_TOP + Cm(0.1),
               img_w, img_h, "ANP Painéis Dinâmicos")

# Painel lateral
panel_l = CONTENT_L + img_w + Cm(0.4)
panel_w = SLIDE_W - panel_l - Cm(0.3)
top = CONTENT_TOP + Cm(0.1)
panel_h = CONTENT_H - Cm(0.2)

add_rect(slide, panel_l, top, panel_w, panel_h, COR_CINZA_CLARO, COR_VERDE_DESTAQUE, Pt(2))
add_rect(slide, panel_l, top, panel_w, Cm(1.0), COR_VERDE_DESTAQUE)
add_textbox(slide, panel_l + Cm(0.2), top + Cm(0.15),
            panel_w - Cm(0.4), Cm(0.7),
            "ANP — Dados Abertos", font_size=11, bold=True,
            color=COR_BRANCO, align=PP_ALIGN.CENTER)

bullets_anp = [
    "Órgão regulador do setor",
    "BI com dados públicos",
    "Disponível para reguladores",
    "Modelos de transparência",
    "Possível com Python + IA",
]
y_b = top + Cm(1.2)
for item in bullets_anp:
    add_textbox(slide, panel_l + Cm(0.2), y_b, panel_w - Cm(0.4), Cm(0.75),
                "✓ " + item, font_size=10, color=COR_CINZA_TEXTO)
    y_b += Cm(0.95)

add_rect(slide, panel_l + Cm(0.2), y_b + Cm(0.3), panel_w - Cm(0.4), Cm(2.0),
         COR_VERDE_DESTAQUE)
add_textbox(slide, panel_l + Cm(0.3), y_b + Cm(0.45), panel_w - Cm(0.6), Cm(1.7),
            "\"Transparência regulatória = dados abertos + BI + IA\"",
            font_size=9, bold=True, color=COR_BRANCO)

print("  Slide 24 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Publicação Web — Ferramentas
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Do Relatório ao Dashboard Público — Ferramentas de Deploy",
                  subtitle_text="Publique seus resultados na web — sem ser programador")

ferramentas_deploy = [
    ("Streamlit", "streamlit.io",
     "Python puro → app interativo\nGráficos, filtros, uploads\nGratuito na nuvem\nIdeal para análises exploratórias",
     COR_VERMELHO, IMG_STREAMLIT),
    ("Vercel", "vercel.com",
     "Deploy de apps em segundos\nSuporte a Python (FastAPI)\nCDN global gratuito\nUsado no Observatório AGENERSA",
     COR_CINZA_TEXTO, None),
    ("Google Data Studio\n(Looker Studio)", "lookerstudio.google.com",
     "Conecta ao Google Sheets\nDashboards sem código\nGratuito com conta Google\nIdeal para dados regulatórios",
     COR_AZUL_CLARO, None),
]

col_w = CONTENT_W / 3 - Cm(0.3)
col_h = CONTENT_H - Cm(0.3)

for i, (nome, url, descr, cor, img_path) in enumerate(ferramentas_deploy):
    left = CONTENT_L + i * (col_w + Cm(0.3))
    top  = CONTENT_TOP + Cm(0.15)
    add_card_tool(slide, left, top, col_w, col_h, nome, url, descr, cor, img_path)

print("  Slide 25 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Demo 4: NotebookLM
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Demo 4 — NotebookLM: A Apostila que Conversa",
                  subtitle_text="notebooklm.google.com — carregue qualquer PDF e converse com ele")

add_demo_banner(slide, CONTENT_TOP)

y = CONTENT_TOP + Cm(1.1)

# Destaque do produto
add_rect(slide, CONTENT_L, y, CONTENT_W, Cm(2.0), COR_CINZA_CLARO, COR_AZUL_CLARO, Pt(2))
add_textbox(slide, CONTENT_L + Cm(0.5), y + Cm(0.3), CONTENT_W - Cm(1.0), Cm(1.4),
            "Carregar a apostila do Prof. Barateiro no NotebookLM → fazer perguntas técnicas em linguagem natural",
            font_size=12, bold=True, color=COR_AZUL_ABAR)

y += Cm(2.3)

# Exemplos de perguntas
perguntas = [
    ("P1", "Qual é a equação de incerteza do Tipo A definida na apostila?",
     "IA responde com fórmula exata + referência à página"),
    ("P2", "Como o método GUM trata correlações entre variáveis?",
     "IA sintetiza o conceito com exemplos da própria apostila"),
    ("P3", "Quais são os passos para calcular incerteza expandida de um medidor de gás?",
     "IA dá resposta passo a passo, citando as seções relevantes"),
    ("P4", "Gere um podcast de 10 min resumindo os conceitos principais",
     "Recurso de áudio do NotebookLM — ideal para revisão"),
]

for num, pergunta, resposta in perguntas:
    add_rect(slide, CONTENT_L, y, Cm(0.8), Cm(1.5), COR_AZUL_ABAR)
    add_textbox(slide, CONTENT_L, y + Cm(0.35), Cm(0.8), Cm(0.8),
                num, font_size=9, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(1.2), y + Cm(0.05), CONTENT_W - Cm(1.3), Cm(0.7),
                pergunta, font_size=11, bold=False, color=COR_AZUL_ABAR)
    add_textbox(slide, Cm(1.2), y + Cm(0.75), CONTENT_W - Cm(1.3), Cm(0.6),
                "↳ " + resposta, font_size=10, color=COR_VERDE_DESTAQUE)
    y += Cm(1.8)

print("  Slide 26 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Materiais e Próximos Passos (alinhamento corrigido)
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Materiais e Próximos Passos",
                  subtitle_text="Tudo que você pode levar desta aula e como continuar aprendendo")

# Dois blocos lado a lado: Materiais | O que Fazer Agora
half_w = CONTENT_W / 2 - Cm(0.3)
top = CONTENT_TOP + Cm(0.2)

# Materiais disponíveis
add_rect(slide, CONTENT_L, top, half_w, Cm(1.0), COR_AZUL_ABAR)
add_textbox(slide, CONTENT_L + Cm(0.2), top + Cm(0.15),
            half_w - Cm(0.4), Cm(0.7),
            "📦  Materiais Disponíveis", font_size=12, bold=True, color=COR_BRANCO)

materiais = [
    "Notebooks Jupyter (7 arquivos) — análise completa do distrito",
    "Script Python do pipeline (FastAPI + Gemini API)",
    "Planilha de dados do distrito (anonimizada)",
    "Relatório Word gerado automaticamente (exemplo)",
    "Artigos científicos organizados por tema (6 pastas)",
    "Esta apresentação (versão com links ativos)",
    "Acesso ao Observatório de Gás: dbagenersa.vercel.app",
]
y = top + Cm(1.1)
for item in materiais:
    add_textbox(slide, CONTENT_L + Cm(0.3), y, half_w - Cm(0.5), Cm(0.7),
                "• " + item, font_size=9.5, color=COR_CINZA_TEXTO)
    y += Cm(0.8)

# O que fazer agora
right_l = CONTENT_L + half_w + Cm(0.3)
add_rect(slide, right_l, top, half_w, Cm(1.0), COR_VERDE_DESTAQUE)
add_textbox(slide, right_l + Cm(0.2), top + Cm(0.15),
            half_w - Cm(0.4), Cm(0.7),
            "🚀  O Que Fazer Agora", font_size=12, bold=True, color=COR_BRANCO)

acoes = [
    ("Esta semana", "Abra o ChatGPT e cole uma planilha sua — veja o que acontece"),
    ("Esta semana", "Configure NotebookLM com a apostila do Barateiro"),
    ("Próximas 2 semanas", "Tente replicar um dos notebooks no seu contexto"),
    ("1 mês", "Implemente um dashboard simples com Streamlit"),
    ("3 meses", "Automatize um relatório recorrente do seu trabalho"),
]
y = top + Cm(1.1)
for prazo, acao in acoes:
    add_rect(slide, right_l, y, Cm(3.5), Cm(0.65), COR_CINZA_MEDIO)
    add_textbox(slide, right_l + Cm(0.1), y + Cm(0.1), Cm(3.3), Cm(0.5),
                prazo, font_size=8, bold=True, color=COR_AZUL_ABAR)
    add_textbox(slide, right_l + Cm(3.7), y + Cm(0.1), half_w - Cm(3.9), Cm(0.55),
                acao, font_size=9.5, color=COR_CINZA_TEXTO)
    y += Cm(0.9)

print("  Slide 27 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Avaliação
# ═══════════════════════════════════════════════════════════════════════════════
slide = std_slide(prs, "Avaliação — Bloco 6",
                  subtitle_text="10 questões de múltipla escolha — 5 minutos")

# Banner central
add_rect(slide, CONTENT_L, CONTENT_TOP + Cm(0.3), CONTENT_W, Cm(3.5), COR_AZUL_ABAR)
add_textbox(slide, CONTENT_L + Cm(1.0), CONTENT_TOP + Cm(0.7), CONTENT_W - Cm(2.0), Cm(2.5),
            "Acesse o formulário de avaliação pelo QR Code\nou pelo link compartilhado no chat",
            font_size=16, bold=True, color=COR_BRANCO, align=PP_ALIGN.CENTER)

# QR Code placeholder
qr_size = Cm(7.0)
qr_l = SLIDE_W / 2 - qr_size / 2
qr_t = CONTENT_TOP + Cm(4.5)
add_rect(slide, qr_l, qr_t, qr_size, qr_size, COR_BRANCO, COR_AZUL_ABAR, Pt(3))
add_textbox(slide, qr_l + Cm(0.5), qr_t + Cm(2.5), qr_size - Cm(1.0), Cm(2.0),
            "[QR Code\nda Avaliação]", font_size=14, color=COR_AZUL_ABAR,
            align=PP_ALIGN.CENTER)

# Temas da avaliação
temas_top = CONTENT_TOP + Cm(4.5)
temas_esq = [
    "Q1 — Capacidades atuais da IA",
    "Q2 — Ferramenta adequada por tarefa",
    "Q3 — Ciclo IBM com IA",
    "Q4 — Web scraping com IA",
    "Q5 — Geração de relatórios",
]
temas_dir = [
    "Q6 — NotebookLM",
    "Q7 — Paralelismo no pipeline",
    "Q8 — Auditoria regulatória",
    "Q9 — Democratização da ciência de dados",
    "Q10 — Integração do ciclo completo",
]

col_w_tema = Cm(8.5)
for i, tema in enumerate(temas_esq):
    add_textbox(slide, CONTENT_L, temas_top + i * Cm(0.9), col_w_tema, Cm(0.8),
                tema, font_size=9.5, color=COR_CINZA_TEXTO)
for i, tema in enumerate(temas_dir):
    right_x = SLIDE_W - MARGIN - col_w_tema
    add_textbox(slide, right_x, temas_top + i * Cm(0.9), col_w_tema, Cm(0.8),
                tema, font_size=9.5, color=COR_CINZA_TEXTO)

print("  Slide 28 OK")


# ═══════════════════════════════════════════════════════════════════════════════
# SALVAR
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT_PPTX)
print(f"\nApresentacao v2 salva: {OUTPUT_PPTX}")
print(f"  Total: {len(prs.slides)} slides")
