"""
Script para gerar a apresentação PowerPoint do Curso ABAR.
Análise de Condições de Operação de Distrito usando Python + AI.
Versão 2: com gráficos dos notebooks incorporados nos slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Cores do tema
AZUL_ESCURO = RGBColor(0x1A, 0x23, 0x7E)
AZUL_MEDIO = RGBColor(0x21, 0x96, 0xF3)
AZUL_CLARO = RGBColor(0xBB, 0xDE, 0xFB)
LARANJA = RGBColor(0xFF, 0x98, 0x00)
VERDE = RGBColor(0x4C, 0xAF, 0x50)
VERMELHO = RGBColor(0xF4, 0x43, 0x36)
ROXO = RGBColor(0x9C, 0x27, 0xB0)
CINZA_ESCURO = RGBColor(0x33, 0x33, 0x33)
CINZA_CLARO = RGBColor(0xF5, 0xF5, 0xF5)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

from config import GRAFICOS_DIR as _GRAFICOS_PATH, PRESENT_DIR
GRAFICOS_DIR = str(_GRAFICOS_PATH)

def add_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=CINZA_ESCURO, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    font_color=CINZA_ESCURO, bullet_char='\u2022'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'{bullet_char} {item}'
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox


def add_code_block(slide, left, top, width, height, code, font_size=12):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x26, 0x32, 0x38)
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    p.font.name = 'Consolas'
    return shape


def add_formula_box(slide, left, top, width, height, formula, font_size=18,
                    bg_color=RGBColor(0xFF, 0xF8, 0xE1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = LARANJA
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(15)
    tf.margin_right = Pt(15)
    p = tf.paragraphs[0]
    p.text = formula
    p.font.size = Pt(font_size)
    p.font.color.rgb = CINZA_ESCURO
    p.font.bold = True
    p.font.name = 'Cambria Math'
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_title_bar(slide, title_text, subtitle_text=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(1.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_ESCURO
    shape.line.fill.background()
    shape.shadow.inherit = False

    add_textbox(slide, 0.5, 0.15, 12, 0.7, title_text,
                font_size=32, font_color=BRANCO, bold=True,
                alignment=PP_ALIGN.LEFT)

    if subtitle_text:
        add_textbox(slide, 0.5, 0.75, 12, 0.5, subtitle_text,
                    font_size=16, font_color=AZUL_CLARO, bold=False,
                    alignment=PP_ALIGN.LEFT)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.3), Inches(13.333), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LARANJA
    line.line.fill.background()
    line.shadow.inherit = False


def add_result_box(slide, left, top, width, height, label, value,
                   color=VERDE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)

    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = BRANCO
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(22)
    p2.font.color.rgb = BRANCO
    p2.font.bold = True
    p2.font.name = 'Calibri'
    p2.alignment = PP_ALIGN.CENTER


def add_image(slide, left, top, width, filename):
    """Adiciona imagem PNG ao slide se o arquivo existir."""
    path = os.path.join(GRAFICOS_DIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))
        return True
    else:
        # Placeholder com borda quando imagem não existe
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(width * 0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CINZA_CLARO
        shape.line.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
        shape.line.width = Pt(1)
        shape.shadow.inherit = False
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f'[Gráfico: {filename}]'
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)
        p.alignment = PP_ALIGN.CENTER
        return False


def gerar_apresentacao():
    """Gera a apresentação PPTX completa. Retorna o caminho do arquivo salvo."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ============================================================
    # SLIDE 1 - CAPA
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, AZUL_ESCURO)

    line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LARANJA
    line.line.fill.background()
    line.shadow.inherit = False

    add_textbox(slide, 1, 1.5, 11.3, 1.2,
            'Análise de Condições de Operação de Distrito',
            font_size=40, font_color=BRANCO, bold=True,
            alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1, 2.8, 11.3, 0.8,
            'Reprodução da Análise com Python e Inteligência Artificial',
            font_size=24, font_color=AZUL_CLARO, bold=False,
            alignment=PP_ALIGN.CENTER)

    div = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(4), Inches(3.8), Inches(5.3), Inches(0.04)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = LARANJA
    div.line.fill.background()
    div.shadow.inherit = False

    add_textbox(slide, 1, 4.2, 11.3, 0.6,
            'Curso ABAR - Medições Inteligentes e Gestão Integrada',
            font_size=22, font_color=LARANJA, bold=True,
            alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1, 5.2, 11.3, 0.5,
            'AGENERSA - Agência Reguladora de Energia e Saneamento do Estado do Rio de Janeiro',
            font_size=16, font_color=RGBColor(0x90, 0xCA, 0xF9), bold=False,
            alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1, 6.2, 11.3, 0.5,
            'Ferramentas: Python | pandas | matplotlib | Claude AI',
            font_size=14, font_color=RGBColor(0x78, 0x78, 0x78), bold=False,
            alignment=PP_ALIGN.CENTER)


    # ============================================================
    # SLIDE 2 - OBJETIVO DO EXERCICIO
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Objetivo do Exercício',
              'O que faremos neste curso prático')

    add_textbox(slide, 0.8, 1.7, 6, 0.5,
            'Reproduzir a análise de dados reais de um distrito de gás usando Python + AI',
            font_size=18, font_color=CINZA_ESCURO, bold=True)

    items_obj = [
    'Ler e explorar dados de uma planilha Excel real de operação',
    'Analisar volumes de entrada (Concessionaria vs Transportadora)',
    'Avaliar o Poder Calorífico Superior (PCS) do gás',
    'Calcular a energia diária (Volume x PCS)',
    'Traçar perfis de consumo dos 7 clientes do distrito',
    'Calcular incertezas de medição combinadas (método RSS)',
    'Realizar o balanço de massa com bandas de incerteza',
    ]
    add_bullet_list(slide, 0.8, 2.4, 6.5, 4, items_obj, font_size=16)

    shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8), Inches(1.7), Inches(4.8), Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = AZUL_MEDIO
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    add_textbox(slide, 8.3, 1.9, 4.2, 0.4,
            'Dados Disponíveis', font_size=18, font_color=AZUL_ESCURO, bold=True)

    dados_items = [
    '1 distrito de distribuição de gás',
    '7 clientes (pontos de saida)',
    '6 meses de dados (abr-set/2025)',
    '183 dias de dados diários de entrada',
    '~4.400 registros horários por cliente',
    'Dados de volume, pressão, temperatura',
    'PCS e energia diários',
    ]
    add_bullet_list(slide, 8.3, 2.5, 4.2, 3.5, dados_items, font_size=14,
                font_color=AZUL_ESCURO)

    add_textbox(slide, 0.8, 6.5, 12, 0.5,
            'Ferramentas: Python 3.x  |  pandas  |  matplotlib  |  openpyxl  |  Jupyter Notebooks  |  Claude AI',
            font_size=13, font_color=RGBColor(0x75, 0x75, 0x75),
            alignment=PP_ALIGN.CENTER)


    # ============================================================
    # SLIDE 3 - ESTRUTURA DOS DADOS
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Etapa 1: Leitura e Exploração dos Dados',
              'Notebook 01 - Entendendo a estrutura da planilha Excel')

    table_data = [
    ['Aba', 'Conteúdo', 'Tipo', 'Registros'],
    ['Vol Entrada Gas', 'Volumes diários de entrada', 'Diario', '183 dias'],
    ['PCS Ent', 'Poder Calorífico Superior', 'Diario', '183 dias'],
    ['Energia Ent', 'Energia = Vol x PCS', 'Diario', '183 dias'],
    ['Cliente #1 a #7', 'Vol, Pressão, Temperatura', 'Horario', '~4.400/cliente'],
    ['Sumario das Medicoes', 'Resumo de condições', 'Consolidado', '-'],
    ['Incertezas', 'Incerteza por medidor', 'Consolidado', '-'],
    ['Balanço', 'Entrada vs Saídas', 'Consolidado', '-'],
    ]

    rows = len(table_data)
    cols = len(table_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.7),
                                Inches(7), Inches(4)).table

    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = table_data[0][j]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = BRANCO
            paragraph.font.name = 'Calibri'
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_ESCURO

    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = table_data[i][j]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = CINZA_ESCURO
                paragraph.font.name = 'Calibri'
            cell.fill.solid()
            cell.fill.fore_color.rgb = CINZA_CLARO if i % 2 == 0 else BRANCO

    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(1.2)
    table.columns[3].width = Inches(1.3)

    add_code_block(slide, 8.3, 1.7, 4.5, 2.5,
               "import pandas as pd\n\n"
               "ARQUIVO = 'Analise de Condicoes\n"
               "           de Operacao de Distrito.xlsx'\n\n"
               "df = pd.read_excel(\n"
               "    ARQUIVO,\n"
               "    sheet_name='Vol Entrada Gas',\n"
               "    header=1,\n"
               "    usecols='B:F'\n"
               ")",
               font_size=12)

    add_textbox(slide, 8.3, 4.5, 4.5, 0.4,
            'Cuidados na leitura:', font_size=14, font_color=VERMELHO, bold=True)

    cuidados = [
    'Nomes de abas com espaços no final',
    'Colunas vazias (deslocamento)',
    'Clientes sem dados (Cliente #4)',
    'header=1 ou header=2 conforme aba',
    ]
    add_bullet_list(slide, 8.3, 5, 4.5, 2, cuidados, font_size=13,
                font_color=CINZA_ESCURO, bullet_char='\u26A0')


    # ============================================================
    # SLIDE 4 - VOLUMES DE ENTRADA (com grafico)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Etapa 2: Análise de Volumes de Entrada',
              'Notebook 02 - Comparando Concessionaria vs Transportadora')

    add_formula_box(slide, 0.8, 1.7, 5.5, 0.7,
                'Diferença (%) = (Vol_Conc - Vol_Transp) / Vol_Conc x 100',
                font_size=16)

    items_vol = [
    'Série temporal de 183 dias (abril a setembro 2025)',
    'Volume médio diário: ~996.000 Nm³/dia',
    'Diferenças entre medições < 0,01%',
    'Dados consistentes entre Concessionária e Transportadora',
    ]
    add_bullet_list(slide, 0.8, 2.6, 5.5, 2, items_vol, font_size=14)

    add_result_box(slide, 0.8, 4.7, 2.5, 1.0,
               'Volume Total', '~183 Mm3', AZUL_MEDIO)
    add_result_box(slide, 3.6, 4.7, 2.5, 1.0,
               'Dif. Conc vs Transp', '< 0,01%', VERDE)

    add_code_block(slide, 0.8, 5.9, 5.5, 1.0,
               "df['Dif_Pct'] = (df['Vol_Conc'] - df['Vol_Transp'])\n"
               "                / df['Vol_Conc'] * 100",
               font_size=11)

    # Grafico
    add_image(slide, 6.8, 1.7, 6.2, 'vol_entrada_serie.png')


    # ============================================================
    # SLIDE 5 - PCS (com grafico)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Etapa 3: Análise do Poder Calorífico Superior (PCS)',
              'Notebook 03 - Qualidade do gás ao longo do tempo')

    shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = VERDE
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    add_textbox(slide, 1.0, 1.8, 5, 0.3,
            'PCS = Poder Calorífico Superior', font_size=16,
            font_color=VERDE, bold=True)
    add_textbox(slide, 1.0, 2.1, 5, 0.7,
            'Energia total por m³ de gás (kcal/m³).\n'
            'Depende da composição (cromatografia).',
            font_size=13, font_color=CINZA_ESCURO)

    items_pcs = [
    'PCS médio: ~9.539 kcal/m³',
    'Variação: 9.168 a 9.786 kcal/m³',
    'Medições Conc vs Transp: dif < 0,01%',
    'Composição varia ao longo do tempo',
    ]
    add_bullet_list(slide, 0.8, 3.1, 5.5, 2, items_pcs, font_size=14)

    add_result_box(slide, 0.8, 5.0, 2.5, 0.9,
               'PCS Médio', '9.539 kcal/m3', VERDE)
    add_result_box(slide, 3.6, 5.0, 1.2, 0.9,
               'Min', '9.168', RGBColor(0x66, 0xBB, 0x6A))
    add_result_box(slide, 5.0, 5.0, 1.2, 0.9,
               'Max', '9.786', RGBColor(0x66, 0xBB, 0x6A))

    # Grafico
    add_image(slide, 6.8, 1.7, 6.2, 'pcs_serie.png')


    # ============================================================
    # SLIDE 6 - ENERGIA (com grafico)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Etapa 4: Cálculo de Energia',
              'Notebook 04 - Energia = Volume x PCS')

    add_formula_box(slide, 0.8, 1.7, 5.5, 0.8,
                'Energia (kcal) = Volume (Nm3) x PCS (kcal/Nm3)\n'
                'Energia (Gcal) = Energia (kcal) / 1.000.000',
                font_size=15)

    items_en = [
    'Energia diária calculada para 183 dias',
    'Validação: valores conferem com planilha (dif ~ 0%)',
    'Energia média diária: ~9.536 Gcal/dia',
    'Relação volume-energia quase linear (r > 0.999)',
    ]
    add_bullet_list(slide, 0.8, 2.8, 5.5, 2, items_en, font_size=14)

    add_result_box(slide, 0.8, 4.7, 2.5, 1.0,
               'Energia Média', '~9.536 Gcal/d', ROXO)
    add_result_box(slide, 3.6, 4.7, 2.5, 1.0,
               'Validação Planilha', 'Dif ~ 0%', VERDE)

    add_code_block(slide, 0.8, 5.9, 5.5, 1.0,
               "df['Energia'] = df['Volume_Nm3d'] * df['PCS']\n"
               "# Verificação: diferença ~ 0% vs planilha",
               font_size=11)

    # Grafico
    add_image(slide, 6.8, 1.7, 6.2, 'energia_serie.png')


    # ============================================================
    # SLIDE 7 - PERFIS DOS CLIENTES (com grafico)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Etapa 5: Perfis de Consumo dos Clientes',
              'Notebook 05 - 7 clientes, 7 padrões de consumo')

    # Tabela compacta
    client_data = [
    ['Cliente', 'Volume', 'Part.'],
    ['Empresa A', '104,1 Mm3', '57,5%'],
    ['Empresa B', '43,7 Mm3', '24,1%'],
    ['Empresa E', '10,2 Mm3', '5,6%'],
    ['Empresa G', '10,1 Mm3', '5,6%'],
    ['Empresa C', '6,8 Mm3', '3,7%'],
    ['Empresa F', '6,0 Mm3', '3,3%'],
    ['Empresa D', '0,09 Mm3', '0,05%'],
    ]

    rows = len(client_data)
    cols = len(client_data[0])
    table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.7),
                                Inches(5.5), Inches(3.5)).table

    for j in range(cols):
        cell = table.cell(0, j)
        cell.text = client_data[0][j]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = BRANCO
            paragraph.font.name = 'Calibri'
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_ESCURO

    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = client_data[i][j]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = CINZA_ESCURO
                paragraph.font.name = 'Calibri'
            cell.fill.solid()
            if i == 1:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
            else:
                cell.fill.fore_color.rgb = CINZA_CLARO if i % 2 == 0 else BRANCO

    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.0)

    insights = [
    'Empresa A domina (57,5%)',
    'Top 2 = 81,6% do volume',
    'Empresa D = 0,05% (GNV)',
    'Fator de carga revela padrão',
    ]
    add_bullet_list(slide, 0.8, 5.4, 5.5, 1.5, insights, font_size=13)

    # Grafico de participacao
    add_image(slide, 6.8, 1.7, 6.2, 'clientes_participacao.png')


    # ============================================================
    # SLIDE 8 - INCERTEZAS (com grafico)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Etapa 6: Cálculo de Incertezas de Medição',
              'Notebook 06 - Método RSS (Root Sum of Squares)')

    add_formula_box(slide, 0.8, 1.7, 5.5, 0.7,
                'U_combinada = sqrt( u1^2 + u2^2 + ... + un^2 )',
                font_size=16)

    add_textbox(slide, 0.8, 2.6, 5.5, 0.3,
            'Incertezas individuais:', font_size=15,
            font_color=AZUL_ESCURO, bold=True)

    inc_items = [
    'Tramo 101: 1,06% | Tramo 501: 1,09%',
    'Empresa A: 1,33% | Empresa B: 1,61%',
    'Empresa C: 1,34% | Empresa E: 3,05%',
    'Empresa F: 1,48% | Empresa G: 2,80%',
    'Empresa D: 3,58%',
    ]
    add_bullet_list(slide, 0.8, 3.0, 5.5, 2, inc_items, font_size=13)

    add_result_box(slide, 0.8, 5.0, 2.5, 1.0,
               'U entrada (RSS)', '1,52%', AZUL_MEDIO)
    add_result_box(slide, 3.6, 5.0, 2.5, 1.0,
               'U saidas (RSS)', '6,19%', LARANJA)

    add_code_block(slide, 0.8, 6.2, 5.5, 0.8,
               "u_entrada = sqrt(0.0106^2 + 0.0109^2) = 1.52%\n"
               "u_saida = sqrt(0.0133^2 + ... + 0.0280^2) = 6.19%",
               font_size=11)

    # Grafico
    add_image(slide, 6.8, 1.7, 6.2, 'incertezas_barras.png')


    # ============================================================
    # SLIDE 9 - BALANCO DE MASSA (com grafico waterfall)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Etapa 7: Balanço de Massa com Incertezas',
              'Notebook 07 - A diferença é aceitável?')

    add_formula_box(slide, 0.8, 1.7, 5.5, 0.6,
                'Dif (%) = (Entrada - Soma_Saidas) / Entrada x 100',
                font_size=15)

    add_textbox(slide, 0.8, 2.5, 5.5, 0.3,
            'Bandas de variação:', font_size=15,
            font_color=AZUL_ESCURO, bold=True)

    add_formula_box(slide, 0.8, 2.9, 5.5, 0.5,
                'V_min = V x (1 - U%)     V_max = V x (1 + U%)',
                font_size=14)

    items_bal = [
    'Se bandas se sobrepõem → diferença explicada pela incerteza',
    'Se NÃO → perda real ou problema no sistema',
    ]
    add_bullet_list(slide, 0.8, 3.6, 5.5, 0.8, items_bal, font_size=13)

    add_textbox(slide, 0.8, 4.5, 5.5, 1.5,
            'Entrada: 182,9 Mm3  [180,1 - 185,7 Mm3]  (+/- 1,52%)\n'
            'Saidas:  180,9 Mm3  [169,7 - 192,1 Mm3]  (+/- 6,19%)\n'
            'Sobreposição: SIM → Balanço ACEITÁVEL',
            font_size=13, font_color=CINZA_ESCURO)

    add_result_box(slide, 0.8, 5.9, 1.7, 0.9, 'Entrada', '182,9 Mm3', AZUL_MEDIO)
    add_result_box(slide, 2.7, 5.9, 1.7, 0.9, 'Saidas', '180,9 Mm3', LARANJA)
    add_result_box(slide, 4.6, 5.9, 1.7, 0.9, 'Diferença', '1,09%', VERDE)

    # Grafico waterfall
    add_image(slide, 6.8, 1.7, 6.2, 'balanco_waterfall.png')


    # ============================================================
    # SLIDE 10 - DASHBOARD RESULTADO (com grafico)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Resultado Final - Dashboard',
              'Visão consolidada do balanço de massa do distrito')

    # Cards de resultado no topo
    add_result_box(slide, 0.8, 1.7, 2.8, 1.3, 'Volume Entrada', '182,9 Mm3', AZUL_MEDIO)
    add_result_box(slide, 3.9, 1.7, 2.8, 1.3, 'Volume Saida', '180,9 Mm3', LARANJA)
    add_result_box(slide, 7.0, 1.7, 2.0, 1.3, 'Diferença', '1,09%', VERDE)
    add_result_box(slide, 9.3, 1.7, 2.0, 1.3, 'U Entrada', '1,52%', AZUL_MEDIO)
    add_result_box(slide, 11.6, 1.7, 1.5, 1.3, 'Resultado', 'OK', VERDE)

    # Grafico de dashboard
    add_image(slide, 0.8, 3.3, 6.5, 'balanco_dashboard.png')

    # Participacao dos clientes
    add_textbox(slide, 7.8, 3.3, 5, 0.4,
            'Participação no Distrito', font_size=16,
            font_color=AZUL_ESCURO, bold=True)

    clientes_pct = [
    'Empresa A:              104,1 Mm3 (57,5%)',
    'Empresa B: 43,7 Mm3 (24,1%)',
    'Empresa E:           10,2 Mm3 (5,6%)',
    'Empresa G:             10,1 Mm3 (5,6%)',
    'Empresa C:    6,8 Mm3 (3,7%)',
    'Empresa F:              6,0 Mm3 (3,3%)',
    'Empresa D (GNV):     0,09 Mm3 (0,05%)',
    ]
    add_bullet_list(slide, 7.8, 3.8, 5, 3, clientes_pct, font_size=13)

    add_textbox(slide, 7.8, 6.5, 5, 0.5,
            'Bandas se sobrepõem → balanço ACEITÁVEL',
            font_size=14, font_color=VERDE, bold=True)


    # ============================================================
    # SLIDE 11 - CONCLUSOES
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Conclusões e Lições Aprendidas',
              'O que aprendemos com este exercício')

    # Citacao destaque
    shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.5), Inches(1.7), Inches(10.3), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    shape.line.color.rgb = LARANJA
    shape.line.width = Pt(3)
    shape.shadow.inherit = False

    add_textbox(slide, 2, 1.85, 9.3, 0.9,
            '"Você fatura o que mede e não o que produz!"',
            font_size=26, font_color=LARANJA, bold=True,
            alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.8, 3.3, 6, 0.4,
            'Conclusões Técnicas:', font_size=18, font_color=AZUL_ESCURO, bold=True)

    conclusoes = [
    'Balanço de massa é essencial para gestão de distritos',
    'Incertezas definem se diferenças são aceitáveis',
    'Método RSS (GUM 2008) é o padrão para combinar incertezas',
    'Diferença de ~1% está dentro das bandas → balanço aceitável',
    'PCS afeta diretamente o valor energético e comercial',
    'Medidores com alta incerteza impactam todo o balanço',
    ]
    add_bullet_list(slide, 0.8, 3.8, 6.5, 3, conclusoes, font_size=15)

    add_textbox(slide, 8, 3.3, 5, 0.4,
            'Sobre o uso de Python + AI:', font_size=18, font_color=ROXO, bold=True)

    ai_items = [
    'Python permite reproduzir análises de forma transparente',
    'Notebooks documentam cada etapa do processo',
    'AI acelera a escrita de código e identifica erros',
    'Validação cruzada com planilha garante confiabilidade',
    'Gráficos automáticos facilitam a interpretação',
    'Código pode ser reutilizado para outros distritos',
    ]
    add_bullet_list(slide, 8, 3.8, 4.8, 3, ai_items, font_size=15,
                font_color=CINZA_ESCURO)


    # ============================================================
    # SLIDE 12 - ESTRUTURA DOS NOTEBOOKS
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRANCO)
    add_title_bar(slide, 'Estrutura dos Notebooks',
              'Roteiro completo para reprodução da análise')

    notebooks = [
    ('01', 'Leitura e Exploração', 'Ler Excel, mapear abas, tipos de dados', AZUL_MEDIO),
    ('02', 'Volumes de Entrada', 'Série temporal, diferenças Conc vs Transp', AZUL_MEDIO),
    ('03', 'Análise do PCS', 'PCS diário, média móvel, distribuição', VERDE),
    ('04', 'Cálculo de Energia', 'Energia = Vol x PCS, verificação', ROXO),
    ('05', 'Perfis dos Clientes', 'Consumo horário, participação, heatmap', LARANJA),
    ('06', 'Incertezas de Medição', 'RSS, limites regulamentares', VERMELHO),
    ('07', 'Balanço de Massa', 'Entrada vs Saídas, bandas, waterfall', RGBColor(0x00, 0x96, 0x88)),
    ]

    y_pos = 1.7
    for num, titulo, desc, cor in notebooks:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(0.8), Inches(0.65)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = cor
        shape.line.fill.background()
        shape.shadow.inherit = False
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.color.rgb = BRANCO
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_textbox(slide, 1.8, y_pos, 4, 0.35,
                    titulo, font_size=16, font_color=CINZA_ESCURO, bold=True)
        add_textbox(slide, 1.8, y_pos + 0.32, 5, 0.3,
                    desc, font_size=12, font_color=RGBColor(0x75, 0x75, 0x75))

        y_pos += 0.75

    add_textbox(slide, 8, 1.7, 5, 0.4,
            'Como executar:', font_size=18, font_color=AZUL_ESCURO, bold=True)

    add_code_block(slide, 8, 2.2, 4.8, 3.5,
               "# 1. Instalar dependencias\n"
               "pip install -r requirements.txt\n\n"
               "# 2. Abrir Jupyter\n"
               "jupyter notebook\n\n"
               "# 3. Executar cada notebook\n"
               "#    em ordem (01 a 07)\n\n"
               "# Arquivo de dados necessario:\n"
               "# 'Analise de Condicoes de\n"
               "#  Operacao de Distrito.xlsx'\n"
               "# (na mesma pasta dos notebooks)",
               font_size=13)

    add_textbox(slide, 8, 6, 4.8, 0.4,
            'Requisitos: Python 3.9+', font_size=13,
            font_color=RGBColor(0x75, 0x75, 0x75))

    add_textbox(slide, 8, 6.3, 4.8, 0.6,
            'pandas, openpyxl, matplotlib, seaborn, numpy',
            font_size=12, font_color=RGBColor(0x75, 0x75, 0x75))


    # ============================================================
    # SLIDE 13 - ENCERRAMENTO
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, AZUL_ESCURO)

    line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LARANJA
    line.line.fill.background()
    line.shadow.inherit = False

    add_textbox(slide, 1, 2.0, 11.3, 1,
            'Obrigado!',
            font_size=48, font_color=BRANCO, bold=True,
            alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1, 3.2, 11.3, 0.8,
            'Curso ABAR - Medições Inteligentes e Gestão Integrada',
            font_size=22, font_color=LARANJA, bold=False,
            alignment=PP_ALIGN.CENTER)

    div = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(4), Inches(4.2), Inches(5.3), Inches(0.04)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = LARANJA
    div.line.fill.background()
    div.shadow.inherit = False

    add_textbox(slide, 1, 4.6, 11.3, 0.8,
            'Análise reproduzida com Python + Claude AI\n'
            'Todos os notebooks e dados estão disponíveis para prática',
            font_size=18, font_color=AZUL_CLARO, bold=False,
            alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1, 5.8, 11.3, 0.5,
            'AGENERSA - Agência Reguladora de Energia e Saneamento do Estado do Rio de Janeiro',
            font_size=14, font_color=RGBColor(0x78, 0x78, 0x78), bold=False,
            alignment=PP_ALIGN.CENTER)


    # ============================================================
    # SALVAR
    # ============================================================
    PRESENT_DIR.mkdir(parents=True, exist_ok=True)
    output_file = str(PRESENT_DIR / 'Apresentacao_Curso_ABAR_Analise_Dados.pptx')
    prs.save(output_file)
    print(f'Apresentacao salva em: {output_file}')
    print(f'Total de slides: {len(prs.slides)}')

    # Verificar graficos incluidos
    graficos_usados = [
    'vol_entrada_serie.png', 'pcs_serie.png', 'energia_serie.png',
    'clientes_participacao.png', 'incertezas_barras.png',
    'balanco_waterfall.png', 'balanco_dashboard.png'
    ]
    for g in graficos_usados:
        path = os.path.join(GRAFICOS_DIR, g)
        status = 'OK' if os.path.exists(path) else 'FALTANDO'
        print(f'  [{status}] {g}')

    return output_file


if __name__ == "__main__":
    gerar_apresentacao()
